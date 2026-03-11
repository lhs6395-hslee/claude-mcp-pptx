# -*- coding: utf-8 -*-
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_shapes(shapes):
    """그룹 내부까지 재귀 탐색하여 모든 도형을 반환"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape

def update_paragraph_text_only(paragraph, new_text):
    """
    [핵심] 문단 객체를 삭제하거나 새로 만들지 않고,
    기존 문단 안에 있는 텍스트(Run)만 교체하여 '줄 간격'과 '스타일'을 보존합니다.
    """
    # 1. 런(Run)이 없으면 최소한 하나 생성
    if not paragraph.runs:
        paragraph.add_run()

    # 2. 첫 번째 런에 텍스트 덮어쓰기
    paragraph.runs[0].text = new_text

    # 3. 뒤따르는 런(잔여 텍스트)만 제거 (문단 자체는 건드리지 않음)
    for i in range(len(paragraph.runs) - 1, 0, -1):
        paragraph._p.remove(paragraph.runs[i]._r)

    # 4. 텍스트가 비어있으면 불렛(점) 제거 처리
    if new_text == "":
        pPr = paragraph._p.get_or_add_pPr()
        buNone = pPr.find(qn('a:buNone'))
        if buNone is None:
            buNone = pPr.makeelement(qn('a:buNone'))
            pPr.insert(0, buNone)
    
    # 5. 줄 간격을 XML 레벨에서 명시적으로 설정 (2.5)
    # PowerPoint가 제대로 인식하도록 XML 직접 조작
    pPr = paragraph._p.get_or_add_pPr()
    
    # 기존 lnSpc 요소 제거
    for lnSpc in pPr.findall(qn('a:lnSpc')):
        pPr.remove(lnSpc)
    
    # 새로운 lnSpc 요소를 pPr의 첫 번째 자식으로 삽입
    # ⚠️ PowerPoint는 pPr 안의 요소 순서에 민감 — lnSpc는 반드시 첫 번째여야 함
    from lxml import etree as _etree
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    lnSpc = _etree.Element(f'{{{A_NS}}}lnSpc')
    spcPct = _etree.SubElement(lnSpc, f'{{{A_NS}}}spcPct')
    spcPct.set('val', '250000')
    pPr.insert(0, lnSpc)

def copy_paragraph_format(source_para, target_para):
    """원본 문단의 서식을 대상 문단에 복사"""
    from lxml import etree as _etree
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # pPr 전체를 deep copy (정렬, 줄간격 등 모두 포함)
    src_pPr = source_para._p.find(f'{{{A_NS}}}pPr')
    tgt_pPr = target_para._p.find(f'{{{A_NS}}}pPr')
    if src_pPr is not None:
        import copy
        new_pPr = copy.deepcopy(src_pPr)
        if tgt_pPr is not None:
            target_para._p.remove(tgt_pPr)
        target_para._p.insert(0, new_pPr)

    # 런이 없으면 먼저 추가
    if not target_para.runs:
        target_para.add_run()

    # 소스 런 속성(rPr) deep copy
    if source_para.runs:
        src_rPr = source_para.runs[0]._r.find(f'{{{A_NS}}}rPr')
        if src_rPr is not None:
            import copy
            tgt_r = target_para.runs[0]._r
            tgt_rPr = tgt_r.find(f'{{{A_NS}}}rPr')
            if tgt_rPr is not None:
                tgt_r.remove(tgt_rPr)
            tgt_r.insert(0, copy.deepcopy(src_rPr))

def update_toc_slide(slide, toc_items, start_number=1):
    """
    목차 슬라이드 업데이트
    - 5개 이하: 단일 페이지
    - 6개 이상: 여러 페이지로 자동 분할 (5개씩)
    
    Args:
        slide: 목차 슬라이드
        toc_items: 목차 항목 리스트
        start_number: 시작 번호 (기본값: 1)
    
    Note: 6개 이상인 경우 이 함수는 첫 5개만 처리하고,
    나머지는 server.py에서 추가 목차 슬라이드를 생성해야 함
    """
    HEADER_LIMIT = Inches(1.8)
    MAX_ITEMS_PER_PAGE = 5
    
    # 5개 초과 시 경고 (server.py가 처리해야 함)
    if len(toc_items) > MAX_ITEMS_PER_PAGE:
        print(f"   ⚠️ 목차 항목이 {len(toc_items)}개입니다. 첫 {MAX_ITEMS_PER_PAGE}개만 표시합니다.")
        print(f"   ⚠️ 나머지 {len(toc_items) - MAX_ITEMS_PER_PAGE}개는 추가 목차 페이지가 필요합니다.")
        toc_items = toc_items[:MAX_ITEMS_PER_PAGE]

    # 1. 텍스트 박스 수집
    candidates = []
    for s in iter_shapes(slide.shapes):
        if not s.has_text_frame: continue
        if s.top < HEADER_LIMIT: continue # 헤더 보호

        # 잡음 제거
        txt = s.text_frame.text.strip()
        if any(x in txt for x in ["GS Neotek", "PAGE", "00/00"]): continue

        candidates.append(s)

    if not candidates:
        print("   ⚠️ 목차 영역을 찾지 못했습니다.")
        return

    # 2. [패턴 인식] "줄이 많은 상자(3줄 이상)" 우선 탐색
    # 템플릿의 '숫자통', '제목통'을 찾습니다.
    multiline_boxes = [s for s in candidates if len(s.text_frame.paragraphs) >= 3]
    multiline_boxes.sort(key=lambda s: s.left) # 좌->우 정렬 (왼쪽:숫자, 오른쪽:제목)

    # -------------------------------------------------------
    # [CASE A] 다중 문단 모드 (기존 줄 간격 유지 필수)
    # -------------------------------------------------------
    if len(multiline_boxes) > 0:
        print(f"   🚀 [다중 문단 모드] 기존 문단 객체를 재활용하여 줄 간격을 유지합니다.")

        # 숫자 박스와 제목 박스가 분리된 구조 (가장 일반적)
        if len(multiline_boxes) >= 2:
            num_box = multiline_boxes[0]
            title_box = multiline_boxes[1]

            # [중요] 데이터가 기존 문단 개수보다 많으면 새 문단 추가
            num_paragraphs = num_box.text_frame.paragraphs
            title_paragraphs = title_box.text_frame.paragraphs
            
            # 필요한 만큼 문단 추가 (템플릿 서식 복사)
            while len(num_paragraphs) < len(toc_items):
                # 마지막 문단의 서식을 복사하여 새 문단 추가
                new_para = num_box.text_frame.add_paragraph()
                copy_paragraph_format(num_paragraphs[-1], new_para)
                num_paragraphs = num_box.text_frame.paragraphs
            
            while len(title_paragraphs) < len(toc_items):
                new_para = title_box.text_frame.add_paragraph()
                copy_paragraph_format(title_paragraphs[-1], new_para)
                title_paragraphs = title_box.text_frame.paragraphs

            # 1. 숫자통 처리
            for i in range(len(num_paragraphs)):
                if i < len(toc_items):
                    # 데이터 채우기 (start_number부터 시작)
                    update_paragraph_text_only(num_paragraphs[i], str(start_number + i))
                else:
                    # 남는 줄 비우기 (지우지 않고 빈칸으로 둠 -> 줄 간격 유지)
                    update_paragraph_text_only(num_paragraphs[i], "")

            # 2. 제목통 처리
            for i in range(len(title_paragraphs)):
                if i < len(toc_items):
                    # 제목 채우기
                    update_paragraph_text_only(title_paragraphs[i], toc_items[i])
                else:
                    # 남는 줄 비우기
                    update_paragraph_text_only(title_paragraphs[i], "")

        # 통짜 박스 하나인 경우
        elif len(multiline_boxes) == 1:
            box = multiline_boxes[0]
            paragraphs = box.text_frame.paragraphs
            for i in range(len(paragraphs)):
                if i < len(toc_items):
                    update_paragraph_text_only(paragraphs[i], toc_items[i])
                else:
                    update_paragraph_text_only(paragraphs[i], "")

    # -------------------------------------------------------
    # [CASE B] 개별 박스 모드 (Fallback)
    # -------------------------------------------------------
    else:
        print("   🚀 [개별 박스 모드] 행 단위로 처리합니다.")
        candidates.sort(key=lambda s: s.top)
        rows = []
        if candidates:
            current = [candidates[0]]
            for i in range(1, len(candidates)):
                if abs(candidates[i].top - candidates[i-1].top) < Inches(0.2):
                    current.append(candidates[i])
                else:
                    rows.append(current)
                    current = [candidates[i]]
            rows.append(current)

        for i, row in enumerate(rows):
            row.sort(key=lambda s: s.left)
            if i < len(toc_items):
                if len(row) >= 2:
                    update_paragraph_text_only(row[0].text_frame.paragraphs[0], str(i+1))
                    update_paragraph_text_only(row[1].text_frame.paragraphs[0], toc_items[i])
                    for extra in row[2:]: extra.text_frame.clear()
                elif len(row) == 1:
                    update_paragraph_text_only(row[0].text_frame.paragraphs[0], toc_items[i])
            else:
                for s in row: s.text_frame.clear()

    print(f"   ✅ 목차 업데이트 완료.")
