# -*- coding: utf-8 -*-
import random
import urllib.request
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# [A] 폰트 & 색상 (가시성 최우선)
FONTS = {
    "HEAD_TITLE": "프리젠테이션 7 Bold",
    "HEAD_DESC":  "프리젠테이션 5 Medium",
    "BODY_TITLE": "Freesentation",
    "BODY_TEXT":  "Freesentation"
}

COLORS = {
    "PRIMARY":    RGBColor(0, 67, 218),    # 파랑 (제목)
    "BLACK":      RGBColor(0, 0, 0),       # 검정 (본문 가시성 강제)
    "DARK_GRAY":  RGBColor(33, 33, 33),    # 진한 회색
    "GRAY":       RGBColor(80, 80, 80),    # 설명글
    "BG_BOX":     RGBColor(248, 249, 250), # 박스 배경
    "BG_WHITE":   RGBColor(255, 255, 255), # 흰색 배경
    "BORDER":     RGBColor(220, 220, 220), # 테두리
    "TERMINAL_BG": RGBColor(48, 10, 36),   # 터미널 배경 (Ubuntu 보라색)
    "TERMINAL_TITLEBAR": RGBColor(44, 44, 44), # 터미널 타이틀 바 (어두운 회색)
    "TERMINAL_TEXT": RGBColor(102, 204, 102),  # 터미널 텍스트 (초록색)
    "TERMINAL_COMMENT": RGBColor(150, 150, 150), # 터미널 주석 (회색)
    # Syntax highlighting
    "SYN_KEY":     RGBColor(102, 217, 239),  # YAML key / 키워드 (하늘색)
    "SYN_STRING":  RGBColor(255, 167, 89),   # 문자열 (주황색)
    "SYN_NUMBER":  RGBColor(255, 215, 0),    # 숫자 (노란색)
    "SYN_KEYWORD": RGBColor(197, 134, 192),  # true/false/null (보라색)
    "SYN_PLAIN":   RGBColor(248, 248, 242),  # 일반 텍스트 (흰색)
    "TERMINAL_RED": RGBColor(255, 95, 86),    # macOS 빨강 버튼
    "TERMINAL_YELLOW": RGBColor(255, 189, 46), # macOS 노랑 버튼
    "TERMINAL_GREEN": RGBColor(39, 201, 63),  # macOS 초록 버튼
    # 의미 기반 색상 (Semantic Colors - KMS PPT 참조)
    "SEM_RED":       RGBColor(185, 28, 28),   # 주의/필수 (제목)
    "SEM_RED_BG":    RGBColor(254, 242, 242), # 주의/필수 (배경)
    "SEM_RED_TEXT":  RGBColor(127, 29, 29),   # 주의/필수 (본문)
    "SEM_YELLOW":    RGBColor(202, 138,  4),   # 주의/경고 (제목) — risk high
    "SEM_ORANGE":    RGBColor(194, 65, 12),   # 경고/핵심 (제목)
    "SEM_ORANGE_BG": RGBColor(255, 247, 237), # 경고/핵심 (배경)
    "SEM_ORANGE_TEXT": RGBColor(154, 52, 18), # 경고/핵심 (본문)
    "SEM_GREEN":     RGBColor(4, 120, 87),    # 긍정/완료 (제목)
    "SEM_GREEN_BG":  RGBColor(236, 253, 245), # 긍정/완료 (배경)
    "SEM_GREEN_TEXT": RGBColor(6, 95, 70),    # 긍정/완료 (본문)
    "SEM_BLUE":      RGBColor(30, 58, 138),   # 참조/조건 (제목)
    "SEM_BLUE_BG":   RGBColor(239, 246, 255), # 참조/조건 (배경)
    "SEM_BLUE_TEXT": RGBColor(30, 64, 175),   # 참조/조건 (본문)
    "CALLOUT_BG":    RGBColor(30, 58, 138),   # 콜아웃 배경 (진한 파랑)
    "CALLOUT_TEXT":  RGBColor(219, 234, 254), # 콜아웃 본문 (밝은 파랑)
}

# [B] 레이아웃 고정 좌표 (템플릿 원본 준수)
LAYOUT = {
    "SLIDE_TITLE_Y": Inches(0.6),      # 헤더 (상단 고정)
    "SLIDE_DESC_Y":  Inches(0.6),      # 설명글 (상단 고정)
    "BODY_START_Y":  Inches(2.0),      # 본문 시작점
    "BODY_LIMIT_Y":  Inches(7.2),      # 본문 한계선
    "MARGIN_X":      Inches(0.5),
    "SLIDE_W":       Inches(13.333)
}

# [C] 라이브러리 로드
try:
    from duckduckgo_search import DDGS
    HAS_SEARCH_LIB = True
except ImportError:
    HAS_SEARCH_LIB = False

def get_image_from_web(query):
    """이미지 검색 (배경 채우기용)"""
    if not HAS_SEARCH_LIB or not query: return None
    try:
        with DDGS() as ddgs:
            r = list(ddgs.images(f"{query} wallpaper minimal business technology", max_results=1))
            if r:
                req = urllib.request.Request(r[0]['image'], headers={'User-Agent': 'Mozilla/5.0'})
                path = f"img_{random.randint(0,9999)}.jpg"
                with urllib.request.urlopen(req, timeout=3) as res, open(path, 'wb') as f: f.write(res.read())
                return path
    except: pass
    return None

def _download_icon(search_term, save_path):
    """아이콘 웹 다운로드 시도 (DDG 이미지 검색 → icons/ 저장)"""
    if not HAS_SEARCH_LIB:
        return False
    try:
        import urllib.request as _ur
        with DDGS() as ddgs:
            results = list(ddgs.images(
                f"{search_term} icon PNG transparent flat",
                max_results=3, type_image="transparent"
            ))
        for r in results:
            img_url = r.get('image', '')
            if not img_url:
                continue
            try:
                req = _ur.Request(img_url, headers={'User-Agent': 'Mozilla/5.0'})
                with _ur.urlopen(req, timeout=4) as res:
                    data = res.read()
                if len(data) < 500:  # 너무 작은 파일 제외
                    continue
                os.makedirs(os.path.dirname(save_path) or "icons", exist_ok=True)
                with open(save_path, 'wb') as f:
                    f.write(data)
                print(f"   ✅ [아이콘 다운로드 완료] '{search_term}' → {save_path}")
                return True
            except Exception:
                continue
    except Exception as e:
        print(f"   ⚠️ [아이콘 다운로드 실패] '{search_term}': {str(e)[:50]}")
    return False

def draw_icon_search(slide, x, y, size, search_term):
    """
    아이콘 로컬 우선 로드 (v4.0 - 로컬 없으면 웹 다운로드 시도)

    전략:
    1. icons/ 폴더에서 로컬 아이콘 파일 우선 검색
    2. 없으면 DDG 이미지 검색으로 다운로드 → icons/ 저장 후 사용
    3. 다운로드도 실패하면 파란색 원형으로 폴백

    파일명 규칙:
    - "upload arrow" → "icons/upload_arrow.png"
    - "server" → "icons/server.png"
    """
    if not search_term:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        oval.fill.solid(); oval.fill.fore_color.rgb = COLORS["PRIMARY"]
        return

    # 로컬 아이콘 파일 검색
    icon_filename = search_term.replace(" ", "_") + ".png"
    icon_path = os.path.join("icons", icon_filename)

    if os.path.exists(icon_path):
        try:
            slide.shapes.add_picture(icon_path, x, y, width=size, height=size)
            print(f"   ✅ [로컬 아이콘 추가됨] '{search_term}' → {icon_path}")
            return
        except Exception as e:
            print(f"   ⚠️ [로컬 아이콘 로드 실패] '{search_term}': {str(e)[:50]}")

    # 로컬 파일 없음 → 웹 다운로드 시도
    print(f"   ⏬ [아이콘 다운로드 시도] '{search_term}'...")
    if _download_icon(search_term, icon_path) and os.path.exists(icon_path):
        try:
            slide.shapes.add_picture(icon_path, x, y, width=size, height=size)
            return
        except Exception:
            pass

    # 다운로드도 실패 → 파란색 원형 폴백
    print(f"   ⚠️ [아이콘 없음] '{search_term}' (파란색 원형 표시)")
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    oval.fill.solid(); oval.fill.fore_color.rgb = COLORS["PRIMARY"]

def clean_body_placeholders(slide):
    """본문 영역(2.0~7.2)만 청소"""
    for shape in list(slide.shapes):
        if shape.top > Inches(2.0) and shape.top < LAYOUT["BODY_LIMIT_Y"]:
            try: shape._element.getparent().remove(shape._element)
            except: pass

def create_content_box(slide, x, y, w, h, title, body, style="gray", search_q=None, compact=False, terminal=False):
    """
    [만능 박스 생성기]
    - 폰트: 제목 16pt / 본문 14pt (가시성 확보)
    - compact=True: grid_2x2용 작은 폰트 (제목 14pt / 본문 12pt)
    - terminal=True: 터미널 스타일 (macOS 터미널 UI)
    - 이미지: 텍스트가 적고 검색어가 있으면 배경 이미지 자동 삽입
    """
    if w < Inches(1.0): w = Inches(1.0)
    if h < Inches(1.0): h = Inches(1.0)

    # 터미널 모드: 별도 함수 호출
    if terminal:
        create_terminal_box(slide, x, y, w, h, title, body, compact=compact)
        return

    # 일반 모드
    bg = COLORS["BG_BOX"] if style=="gray" else COLORS["BG_WHITE"]
    line = COLORS["BORDER"] if style=="gray" else COLORS["PRIMARY"]

    # 이미지 자동 채우기 비활성화 (안정성 우선)
    filled_image = False
    # 웹 이미지 다운로드는 레이트 리미트 문제로 비활성화

    # 박스 생성
    if not filled_image:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
        text_shape = shp
    else:
        text_shape = slide.shapes.add_textbox(x, y, w, h)

    # 텍스트 설정 (compact 모드에 따라 마진 조정)
    tf = text_shape.text_frame; tf.clear()
    if compact:
        tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.8)  # 오른쪽 여백 증가 (아이콘 공간)
        tf.margin_top = Inches(0.4); tf.margin_bottom = Inches(0.4)  # 위/아래 여백 증가
    else:
        tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.4); tf.margin_bottom = Inches(0.4)  # 위/아래 여백 증가
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 텍스트 오버플로우 방지 (Shrink text on overflow)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # 수직 중앙 정렬

    # 색상 및 폰트
    title_color = COLORS["BG_WHITE"] if filled_image else COLORS["PRIMARY"]
    body_color  = COLORS["BG_WHITE"] if filled_image else COLORS["BLACK"]
    title_font = FONTS["BODY_TITLE"]
    body_font = FONTS["BODY_TEXT"]

    # compact 모드에 따라 폰트 크기 조정
    title_size = Pt(15) if compact else Pt(16)
    body_size = Pt(13) if compact else Pt(14)
    line_spacing = Pt(6) if compact else Pt(8)

    if title:
        p = tf.paragraphs[0]; p.text = str(title)
        p.font.name = title_font; p.font.bold = True; p.font.size = title_size
        p.font.color.rgb = title_color
        p.space_after = line_spacing

    if body:
        # 줄바꿈(\n) 처리를 위해 각 줄을 별도 paragraph로 추가
        lines = str(body).split('\n')
        # 여러 줄이면 개조식(bullet list)으로 간주 → • 자동 추가
        is_list = len(lines) > 1
        for i, line in enumerate(lines):
            if i == 0 and not title:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            stripped = line.strip()
            # 개조식 자동 • 추가: 다중 줄이고, 이미 기호/번호로 시작하지 않는 경우
            # 번호 목록 패턴: "1. ", "2) ", "3: " — 실제 번호 목록만 제외
            import re as _re
            is_numbered = bool(_re.match(r'^\d+[.):\s]\s', stripped))
            if is_list and stripped and not stripped.startswith('•') and not is_numbered:
                p.text = "• " + line
            else:
                p.text = line
            p.font.name = body_font; p.font.size = body_size
            p.font.color.rgb = body_color
            p.alignment = PP_ALIGN.LEFT
            p.space_after = line_spacing

    # 아이콘 추가 (오른쪽 상단, 박스 크기 충분할 때만)
    if search_q:
        icon_size = Inches(0.6)
        # 박스가 아이콘을 수용할 수 있을 때만 표시 (최소 높이 1.2", 너비 1.5")
        if h >= Inches(1.2) and w >= Inches(1.5):
            icon_x = x + w - icon_size - Inches(0.25)
            icon_y_pos = y + Inches(0.2)
            draw_icon_search(slide, icon_x, icon_y_pos, icon_size, search_q)

def _detect_lang(code: str, filename: str = "") -> str:
    """코드 내용과 파일명으로 언어 자동 감지. 명시적 lang이 없을 때 사용."""
    import re as _re
    # 파일명 확장자 우선
    ext_map = {
        '.yaml': 'yaml', '.yml': 'yaml',
        '.py': 'python', '.sh': 'bash', '.bash': 'bash',
        '.tf': 'bash',  # terraform은 bash 스타일로 근사
    }
    for ext, lang in ext_map.items():
        if filename.lower().endswith(ext):
            return lang

    # 코드 내용 휴리스틱
    lines = [l.strip() for l in code.strip().splitlines() if l.strip()][:10]
    sample = '\n'.join(lines)

    yaml_score = sum([
        bool(_re.search(r'^\w[\w\-]*\s*:', sample, _re.M)),   # key: value 패턴
        'apiVersion' in sample,
        'kind:' in sample,
        sample.count(':') > sample.count('='),
    ])
    python_score = sum([
        bool(_re.search(r'^(def |class |import |from )', sample, _re.M)),
        'self.' in sample,
        bool(_re.search(r':\s*$', sample, _re.M)),
    ])
    bash_score = sum([
        bool(_re.search(r'^(kubectl|helm|docker|aws|terraform|git|echo|\$)', sample, _re.M)),
        sample.count('$') > 1,
        bool(_re.search(r'^#!/', sample)),
    ])

    scores = {'yaml': yaml_score, 'python': python_score, 'bash': bash_score}
    best = max(scores, key=scores.get)
    return best if scores[best] >= 2 else None


def _syntax_color(line, lang):
    """줄 단위 syntax highlighting - (text, RGBColor) 튜플 리스트 반환"""
    import re as _re
    stripped = line.strip()
    if stripped.startswith('#'):
        return [(line, COLORS["TERMINAL_COMMENT"])]
    if lang == 'yaml':
        m = _re.match(r'^(\s*-?\s*)(\w[\w\-]*)(\s*:)(.*)', line)
        if m:
            indent, key, colon, rest = m.groups()
            parts = [(indent + key, COLORS["SYN_KEY"]), (colon, COLORS["SYN_PLAIN"])]
            rv = rest.strip()
            if rv in ('true', 'false', 'null', 'True', 'False', 'None'):
                parts.append((' ' + rv, COLORS["SYN_KEYWORD"]))
            elif rv and rv[0] in ('"', "'"):
                parts.append((' ' + rv, COLORS["SYN_STRING"]))
            elif _re.match(r'^[\d\.]+$', rv):
                parts.append((' ' + rv, COLORS["SYN_NUMBER"]))
            elif rv:
                parts.append((' ' + rv, COLORS["SYN_PLAIN"]))
            return parts
        m2 = _re.match(r'^(\s*)(-)(\s+)(.*)', line)
        if m2:
            ind, dash, sp, val = m2.groups()
            return [(ind + dash + sp, COLORS["SYN_PLAIN"]), (val, COLORS["SYN_STRING"])]
    elif lang in ('python', 'py'):
        keywords = {'def', 'class', 'import', 'from', 'return', 'if', 'else', 'elif',
                    'for', 'while', 'in', 'not', 'and', 'or', 'True', 'False', 'None',
                    'with', 'as', 'try', 'except', 'raise', 'pass', 'lambda'}
        if _re.search(r'["\']', line):
            return [(line, COLORS["SYN_STRING"])]
        first_word = stripped.split()[0] if stripped.split() else ''
        if first_word in keywords:
            return [(line, COLORS["SYN_KEYWORD"])]
    elif lang in ('bash', 'sh', 'shell'):
        cmd_colors = {'kubectl', 'helm', 'docker', 'aws', 'terraform', 'git', 'curl', 'echo'}
        first_word = stripped.split()[0] if stripped.split() else ''
        if first_word in cmd_colors:
            return [(line, COLORS["SYN_KEY"])]
        if stripped.startswith('$') or '=$' in line:
            return [(line, COLORS["SYN_KEY"])]
    return [(line, COLORS["SYN_PLAIN"])]


def create_terminal_box(slide, x, y, w, h, title, body, compact=False, lang=None):
    """Ubuntu 스타일 터미널 박스 생성 (사각형)"""
    titlebar_h = Inches(0.3)

    # 1. 전체 배경 박스 (Ubuntu 보라색, 사각형)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["TERMINAL_BG"]
    background.line.color.rgb = COLORS["TERMINAL_BG"]

    # 2. 타이틀 바 (어두운 회색)
    titlebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, titlebar_h)
    titlebar.fill.solid()
    titlebar.fill.fore_color.rgb = COLORS["TERMINAL_TITLEBAR"]
    titlebar.line.color.rgb = COLORS["TERMINAL_TITLEBAR"]

    # 3. macOS 버튼 3개
    btn_size = Inches(0.11)
    btn_y = y + (titlebar_h - btn_size) / 2
    btn_colors = [COLORS["TERMINAL_RED"], COLORS["TERMINAL_YELLOW"], COLORS["TERMINAL_GREEN"]]
    btn_gap = Inches(0.06)

    for i, color in enumerate(btn_colors):
        btn_x = x + Inches(0.12) + i * (btn_size + btn_gap)
        btn = slide.shapes.add_shape(MSO_SHAPE.OVAL, btn_x, btn_y, btn_size, btn_size)
        btn.fill.solid()
        btn.fill.fore_color.rgb = color
        btn.line.color.rgb = color

    # 4. "bash" 타이틀
    title_tb = slide.shapes.add_textbox(x + Inches(0.5), y, w - Inches(0.5), titlebar_h)
    tf_title = title_tb.text_frame
    tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_title.margin_left = Inches(0)
    p_title = tf_title.paragraphs[0]
    p_title.text = title if title else "bash"
    p_title.font.name = "Courier New"
    p_title.font.size = Pt(10)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(200, 200, 200)
    p_title.alignment = PP_ALIGN.CENTER

    # 5. 코드 텍스트 (여백 충분히 확보)
    text_y = y + titlebar_h + Inches(0.15)
    text_h = h - titlebar_h - Inches(0.3)

    text_tb = slide.shapes.add_textbox(x + Inches(0.25), text_y, w - Inches(0.5), text_h)
    tf = text_tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 텍스트 오버플로우 방지
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # 수직 중앙 정렬
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    lines = str(body).split('\n')
    n_lines = len(lines)

    # 라인 수에 따라 동적 폰트 크기 조절
    # generate.py에서 split_text_code는 MAX_LINES_PER_SLIDE=14 이하로 분할 전달됨
    if compact or n_lines > 20:
        font_size = Pt(9); line_spacing = Pt(2)
    elif n_lines > 15:
        font_size = Pt(10); line_spacing = Pt(3)
    elif n_lines > 10:
        font_size = Pt(11); line_spacing = Pt(4)
    else:
        font_size = Pt(14); line_spacing = Pt(6)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # syntax highlighting 적용
        segments = _syntax_color(line, lang) if lang else [(line, COLORS["TERMINAL_TEXT"])]

        # 단일 세그먼트면 p.text로 간단히 처리
        if len(segments) == 1:
            p.text = segments[0][0]
            p.font.name = FONTS["BODY_TEXT"]
            p.font.size = font_size
            p.font.color.rgb = segments[0][1]
        else:
            # 다중 세그먼트: run으로 분리
            p.text = ''
            for seg_text, seg_color in segments:
                if not seg_text:
                    continue
                run = p.add_run()
                run.text = seg_text
                run.font.name = FONTS["BODY_TEXT"]
                run.font.size = font_size
                run.font.color.rgb = seg_color

        p.alignment = PP_ALIGN.LEFT
        p.space_after = line_spacing

def set_slide_title_area(slide, title_text, desc_text=""):
    """헤더 설정 (템플릿 좌표 준수)"""
    # 1. 제목 (Left)
    title_shape = slide.shapes.title or slide.shapes.add_textbox(Inches(0.5), LAYOUT["SLIDE_TITLE_Y"], Inches(4.5), Inches(1.0))
    title_shape.left, title_shape.top = Inches(0.5), LAYOUT["SLIDE_TITLE_Y"]
    title_shape.width = Inches(4.5)

    tf = title_shape.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    
    # 항상 28pt 사용 (CLAUDE.md 규칙: 폰트 크기를 줄이지 말 것)
    # 제목이 길면 \n으로 개행하여 처리
    _title_pt = Pt(28)
    
    # \n 개행 지원: 줄별로 별도 paragraph 생성
    _title_parts = str(title_text).split('\n')
    p.text = _title_parts[0]
    p.font.name = FONTS["HEAD_TITLE"]; p.font.size = _title_pt; p.font.bold = True
    p.font.color.rgb = COLORS["PRIMARY"]; p.alignment = PP_ALIGN.LEFT
    for _part in _title_parts[1:]:
        _p2 = tf.add_paragraph()
        _p2.text = _part
        _p2.font.name = FONTS["HEAD_TITLE"]; _p2.font.size = _title_pt; _p2.font.bold = True
        _p2.font.color.rgb = COLORS["PRIMARY"]; _p2.alignment = PP_ALIGN.LEFT

    # 2. 설명 (Right)
    desc_box = None
    for s in slide.shapes:
        if s.has_text_frame and s.left > Inches(5.0) and s.top < Inches(1.5):
            desc_box = s; break
    if not desc_box:
        desc_box = slide.shapes.add_textbox(Inches(5.2), LAYOUT["SLIDE_DESC_Y"], Inches(7.6), Inches(1.2))

    tf_d = desc_box.text_frame; tf_d.clear(); tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]; p_d.text = str(desc_text)
    p_d.font.name = FONTS["HEAD_DESC"]; p_d.font.size = Pt(12); p_d.font.color.rgb = COLORS["GRAY"]
    p_d.alignment = PP_ALIGN.LEFT

def draw_body_header_and_get_y(slide, title, desc):
    """본문 헤더 (동적 위치 계산)"""
    current_y = LAYOUT["BODY_START_Y"]
    content_w = LAYOUT["SLIDE_W"] - (LAYOUT["MARGIN_X"] * 2)

    if title:
        tb = slide.shapes.add_textbox(LAYOUT["MARGIN_X"], current_y, content_w, Inches(0.6))
        title_str = str(title)
        p = tb.text_frame.paragraphs[0]; p.text = title_str if title_str.startswith("•") else "• " + title_str
        p.font.name = FONTS["BODY_TITLE"]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLORS["DARK_GRAY"]
        current_y += Inches(0.6)  # 박스 높이와 일치

        if desc:
            tb_d = slide.shapes.add_textbox(LAYOUT["MARGIN_X"], current_y, content_w, Inches(0.5))
            tb_d.text_frame.word_wrap = True
            p_d = tb_d.text_frame.paragraphs[0]; p_d.text = str(desc)
            p_d.font.name = FONTS["BODY_TEXT"]; p_d.font.size = Pt(12); p_d.font.color.rgb = COLORS["GRAY"]
            current_y += Inches(0.55)  # 박스 높이 + 소폭 여백
        current_y += Inches(0.15)
    return current_y

def calculate_dynamic_rect(start_y):
    """남은 공간 계산"""
    available_h = LAYOUT["BODY_LIMIT_Y"] - start_y
    if available_h < Inches(1.5): available_h = Inches(1.5)
    return LAYOUT["MARGIN_X"], start_y, LAYOUT["SLIDE_W"] - (LAYOUT["MARGIN_X"] * 2), available_h

