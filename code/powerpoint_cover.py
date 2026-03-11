# -*- coding: utf-8 -*-
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR # 테마 색상 인식용

def replace_text_preserving_style(shape, new_text):
    """
    [헬퍼] 텍스트 교체 (첫 문단 스타일 유지, 나머지 삭제)
    """
    if not shape.has_text_frame: return
    tf = shape.text_frame
    tf.word_wrap = True # 줄바꿈 허용

    if not tf.paragraphs:
        tf.text = new_text; return

    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.text = new_text
    else:
        p0.runs[0].text = new_text
        for i in range(len(p0.runs) - 1, 0, -1):
            p0._p.remove(p0.runs[i]._r)

    # 첫 문단 이후의 나머지 문단(잔존 텍스트) 모두 삭제
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p_element = tf.paragraphs[i]._p
        p_element.getparent().remove(p_element)

def find_shapes_by_keywords(shapes, keywords):
    """키워드 포함 도형 검색 (재귀)"""
    found = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found.extend(find_shapes_by_keywords(shape.shapes, keywords))
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            for kw in keywords:
                if kw in text:
                    found.append(shape)
                    break
    return found

def center_shape_horizontally(shape, slide_width_inch=13.333, shape_width_inch=None):
    """
    [NEW] 도형을 지정된 너비로 설정하고, 슬라이드 정중앙에 배치하는 함수
    """
    if shape_width_inch:
        shape.width = Inches(shape_width_inch)

    # 중앙 좌표 계산: (슬라이드너비 - 도형너비) / 2
    slide_width = Inches(slide_width_inch)
    shape.left = int((slide_width - shape.width) / 2)

    # 텍스트 내부 중앙 정렬
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

# [핵심 1] 템플릿의 색상 정보(RGB 또는 테마 색상) 추출
def get_original_style(shape):
    style = {
        'name': None, 'size': None, 'bold': None,
        'italic': None, 'color_rgb': None, 'color_theme': None, 'brightness': 0
    }

    if shape.has_text_frame and shape.text_frame.paragraphs:
        try:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                style['name'] = r.font.name
                style['size'] = r.font.size
                style['bold'] = r.font.bold
                style['italic'] = r.font.italic

                if r.font.color.type == MSO_COLOR_TYPE.RGB:
                    style['color_rgb'] = r.font.color.rgb
                elif r.font.color.type == MSO_COLOR_TYPE.THEME:
                    style['color_theme'] = r.font.color.theme_color
                    style['brightness'] = r.font.color.brightness
        except: pass
    return style

# [핵심 2] 스타일 승계 및 텍스트 교체
def apply_text_with_style(shape, text, inherited_style, force_center=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = tf.paragraphs[0]

    if force_center:
        p.alignment = PP_ALIGN.CENTER
    else:
        p.alignment = inherited_style.get('alignment', PP_ALIGN.CENTER)

    lines = text.split('\\n') if '\\n' in text else text.split('\n')

    for i, line in enumerate(lines):
        run = p.add_run()
        run.text = line

        # 스타일 복원
        if inherited_style['name']: run.font.name = inherited_style['name']
        if inherited_style['size']: run.font.size = inherited_style['size']
        if inherited_style['bold'] is not None: run.font.bold = inherited_style['bold']

        # 색상 복원
        if inherited_style['color_rgb']:
            run.font.color.rgb = inherited_style['color_rgb']
        elif inherited_style['color_theme']:
            run.font.color.theme_color = inherited_style['color_theme']
            if inherited_style['brightness']:
                run.font.color.brightness = inherited_style['brightness']
        else:
            run.font.color.rgb = RGBColor(255, 255, 255)

        if i < len(lines) - 1:
            run.text += '\n'

def center_shape_horizontally(shape, slide_width_inch=13.333, fixed_width_inch=None):
    if fixed_width_inch:
        shape.width = Inches(fixed_width_inch)
    shape.left = int((Inches(slide_width_inch) - shape.width) / 2)

def update_cover_slide(slide, title_text, subtitle_text):
    now = datetime.now()
    current_year = str(now.year)
    current_md = now.strftime("%m/%d")

    title_candidates = []
    subtitle_candidates = []

    # 1. 도형 분류
    for shape in list(slide.shapes):
        if not shape.has_text_frame: continue
        txt = shape.text_frame.text

        # (A) 날짜
        if any(k in txt for k in ["2025", "2026", "02/06", "02.06", "00/00"]):
            style = get_original_style(shape)
            new_text = current_year if ("20" in txt) else current_md
            apply_text_with_style(shape, new_text, style, force_center=False)
            continue

        # (B) 부제목
        if any(k in txt for k in ["설계", "원칙", "부제목", "Subtitle", "소제목"]):
            subtitle_candidates.append(shape)
            continue

        # (C) 제목
        if any(k in txt for k in ["가이드라인", "GS", "Template", "제목", "AWS"]):
            title_candidates.append(shape)

    # 2. 중복 제거 및 스타일 적용
    target_title = None
    if title_candidates:
        target_title = title_candidates[0]
        saved_style = get_original_style(target_title)
        for trash in title_candidates[1:]:
            try: trash._element.getparent().remove(trash._element)
            except: pass
        apply_text_with_style(target_title, title_text, saved_style, force_center=True)
        center_shape_horizontally(target_title, fixed_width_inch=10.5)

    target_subtitle = None
    if subtitle_candidates:
        target_subtitle = subtitle_candidates[0]
        saved_style = get_original_style(target_subtitle)
        for trash in subtitle_candidates[1:]:
            try: trash._element.getparent().remove(trash._element)
            except: pass
        apply_text_with_style(target_subtitle, subtitle_text, saved_style, force_center=True)
        center_shape_horizontally(target_subtitle, fixed_width_inch=11.333)

    # 3. [FIXED] 수직 정렬 계산 (int 변환 추가)
    slide_height = Inches(7.5)
    gap = Inches(0.3)

    if target_title and target_subtitle:
        total_block_height = target_title.height + gap + target_subtitle.height

        # [수정됨] 나눗셈 결과를 int()로 감싸서 정수로 변환
        start_top = int((slide_height - total_block_height) / 2)

        target_title.top = start_top
        target_subtitle.top = target_title.top + target_title.height + gap

        print(f"✅ 표지 완료: 수직 중앙 정렬 적용 (Top: {target_title.top/914400:.2f}in)")

    elif target_title:
        # 제목만 있을 때도 int() 변환 필요
        target_title.top = int((slide_height - target_title.height) / 2)
