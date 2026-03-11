# -*- coding: utf-8 -*-
from powerpoint_utils import *

# 32. Fishbone Cause-Effect (피쉬본 원인-결과)
def render_fishbone_cause_effect(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    effect = content.get('effect', ''); categories = content.get('categories', [])
    if not categories: return
    n = len(categories)

    SLIDE_BOTTOM = Inches(7.4)  # 슬라이드 하단 경계

    # 가용 높이 기반으로 카드/branch 크기 동적 계산
    # 위: by ~ spine_y, 아래: spine_y ~ min(by+bh, SLIDE_BOTTOM)
    # spine을 body 중앙에 배치하되 상하 여백 확보
    half_h = min(int(bh * 0.42), int((SLIDE_BOTTOM - by) * 0.42))
    spine_y = int(by) + int(bh * 0.50)
    # spine이 슬라이드 밖으로 나가지 않도록 클램프
    spine_y = min(spine_y, int(SLIDE_BOTTOM) - half_h - int(Inches(0.1)))

    # 카드 높이: branch 공간의 85% (branch_h = half_h - 커넥터 여백)
    branch_h = half_h - int(Inches(0.15))
    card_h = int(branch_h * 0.90)
    # 카드 상단이 by 위로 올라가지 않도록 branch_h 재조정
    max_branch_up = spine_y - int(by) - int(Inches(0.05))
    if branch_h > max_branch_up:
        branch_h = max_branch_up
        card_h = int(branch_h * 0.90)
    # 카드 하단이 SLIDE_BOTTOM 밖으로 나가지 않도록
    max_branch_down = int(SLIDE_BOTTOM) - spine_y - int(Inches(0.05))
    if branch_h > max_branch_down:
        branch_h = max_branch_down
        card_h = int(branch_h * 0.90)

    spine_x1 = int(bx) + int(Inches(0.2)); spine_x2 = int(bx + bw) - int(Inches(2.2))
    spine = slide.shapes.add_connector(1, spine_x1, spine_y, spine_x2, spine_y)
    spine.line.color.rgb = COLORS["PRIMARY"]; spine.line.width = Pt(3.0)

    # 효과 박스 (오른쪽 화살촉)
    eff_w = Inches(2.2); eff_h = Inches(0.9)
    eff_shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, spine_x2 - int(Inches(0.2)), spine_y - int(eff_h) // 2, int(eff_w), int(eff_h))
    eff_shp.fill.solid(); eff_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; eff_shp.line.color.rgb = COLORS["PRIMARY"]
    tf = eff_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]; p.text = effect; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    usable_w = spine_x2 - spine_x1 - int(Inches(0.3)); spacing = int(usable_w) / max(n, 1)
    default_colors = ['blue', 'green', 'orange', 'red', 'primary', 'gray']

    for i, cat in enumerate(categories):
        is_top = (i % 2 == 0)
        cx = spine_x1 + int(Inches(0.3)) + int(spacing * i) + int(spacing * 0.5)
        color_key = cat.get('color', default_colors[i % len(default_colors)])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])
        end_y = spine_y - branch_h if is_top else spine_y + branch_h

        conn = slide.shapes.add_connector(1, cx, spine_y, cx, end_y)
        conn.line.color.rgb = line_c; conn.line.width = Pt(2.0)

        card_w = min(int(Inches(2.2)), int(spacing * 0.90))
        card_x = cx - card_w // 2
        card_y = (end_y - card_h) if is_top else end_y
        # 경계 클램프
        card_y = max(int(by), card_y)
        card_y = min(int(SLIDE_BOTTOM) - card_h, card_y)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = fill_c; card.line.color.rgb = line_c; card.line.width = Pt(1.5)
        tf = card.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08); tf.margin_top = Inches(0.06)
        p = tf.paragraphs[0]; p.text = cat.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(12); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(3)
        for cause in cat.get('causes', [])[:4]:
            p2 = tf.add_paragraph(); p2.text = f"• {cause}"; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(9); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.LEFT; p2.space_after = Pt(1)


# 33. Org Chart (조직도/트리 계층)
def render_org_chart(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    root = content.get('root', {}); children = content.get('children', [])
    if not root: return

    # 루트 노드
    root_w = Inches(2.5); root_h = Inches(0.8)
    root_x = int(bx) + int(bw) // 2 - int(root_w) // 2; root_y = int(by)
    root_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, root_x, root_y, int(root_w), int(root_h))
    root_shp.fill.solid(); root_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; root_shp.line.color.rgb = COLORS["PRIMARY"]
    tf = root_shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = root.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(18); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER
    if root.get('desc'):
        p2 = tf.add_paragraph(); p2.text = root['desc']; p2.font.name = FONTS["BODY_TEXT"]
        p2.font.size = Pt(12); p2.font.color.rgb = COLORS["BG_WHITE"]; p2.alignment = PP_ALIGN.CENTER

    if not children: return
    n = len(children)

    # 자식 노드 배치 (전체 폭 사용)
    connector_gap = Inches(0.5)
    child_y = int(by) + int(root_h) + int(connector_gap)
    gap_between = Inches(0.2)
    child_w = (int(bw) - int(gap_between) * (n - 1)) // max(n, 1)
    child_h = int(bh) - int(root_h) - int(connector_gap)
    total_w = child_w * n + int(gap_between) * (n - 1)
    start_x = int(bx) + (int(bw) - total_w) // 2
    root_bottom_y = root_y + int(root_h)
    mid_y = root_bottom_y + int(connector_gap) // 2
    root_cx = root_x + int(root_w) // 2

    # 수직선 + 수평 연결선
    vert = slide.shapes.add_connector(1, root_cx, root_bottom_y, root_cx, mid_y)
    vert.line.color.rgb = COLORS["PRIMARY"]; vert.line.width = Pt(2.0)
    first_cx = start_x + child_w // 2; last_cx = start_x + total_w - child_w // 2
    if n > 1:
        horiz = slide.shapes.add_connector(1, first_cx, mid_y, last_cx, mid_y)
        horiz.line.color.rgb = COLORS["PRIMARY"]; horiz.line.width = Pt(2.0)

    # 자식 수에 따라 폰트 크기 동적 조절
    if n >= 5:
        label_sz = Pt(12); desc_sz = Pt(10); item_sz = Pt(9)
    elif n >= 4:
        label_sz = Pt(14); desc_sz = Pt(11); item_sz = Pt(10)
    else:
        label_sz = Pt(16); desc_sz = Pt(12); item_sz = Pt(11)

    default_colors = ['blue', 'green', 'orange', 'red', 'primary', 'gray']
    for i, child in enumerate(children):
        cx = start_x + int((child_w + int(gap_between)) * i)
        ccx = cx + child_w // 2
        cv = slide.shapes.add_connector(1, ccx, mid_y, ccx, child_y)
        cv.line.color.rgb = COLORS["PRIMARY"]; cv.line.width = Pt(2.0)
        color_key = child.get('color', default_colors[i % len(default_colors)])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, child_y, child_w, child_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c; shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]; p.text = child.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = label_sz; p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
        desc = child.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = desc_sz; p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(6)
        for item in child.get('items', [])[:3]:
            p3 = tf.add_paragraph(); p3.text = f"• {item}"; p3.font.name = FONTS["BODY_TEXT"]
            p3.font.size = item_sz; p3.font.color.rgb = text_c; p3.alignment = PP_ALIGN.LEFT; p3.space_after = Pt(2)


# 34. Temple Pillars (기둥형 구조도)
def render_temple_pillars(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    roof = content.get('roof', {}); pillars = content.get('pillars', []); foundation = content.get('foundation', {})
    if not pillars: return
    n = len(pillars)

    roof_h = Inches(0.7); found_h = Inches(0.6); gap = Inches(0.12)
    pillar_h = int(bh) - int(roof_h) - int(found_h) - int(gap) * 2

    # 지붕 (삼각 페디먼트)
    roof_shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, int(bx), int(by), int(bw), int(roof_h))
    roof_shp.fill.solid(); roof_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; roof_shp.line.color.rgb = COLORS["PRIMARY"]
    tf = roof_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = roof.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(14); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    # 기둥들
    pillar_y = int(by) + int(roof_h) + int(gap)
    pillar_gap = Inches(0.15)
    pillar_w = (int(bw) - int(pillar_gap) * (n - 1)) / max(n, 1)
    default_colors = ['blue', 'green', 'orange', 'red', 'primary', 'gray']
    for i, pil in enumerate(pillars):
        px = int(bx) + int((pillar_w + int(pillar_gap)) * i)
        color_key = pil.get('color', default_colors[i % len(default_colors)])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(px), pillar_y, int(pillar_w), pillar_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c; shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        p = tf.paragraphs[0]; p.text = pil.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(14); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
        desc = pil.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(10); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER

    # 기초 (토대)
    found_y = pillar_y + pillar_h + int(gap)
    found_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(bx), found_y, int(bw), int(found_h))
    found_shp.fill.solid(); found_shp.fill.fore_color.rgb = COLORS["SEM_BLUE_BG"]
    found_shp.line.color.rgb = COLORS["PRIMARY"]; found_shp.line.width = Pt(2.0)
    tf = found_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = foundation.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(13); p.font.color.rgb = COLORS["PRIMARY"]; p.alignment = PP_ALIGN.CENTER


# 35. Infinity Loop (무한 순환 루프)
def render_infinity_loop(slide, data):
    import math
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    left_items = content.get('left_loop', []); right_items = content.get('right_loop', [])
    center_label = content.get('center_label', '')
    left_label = content.get('left_label', ''); right_label = content.get('right_label', '')
    if not left_items and not right_items: return

    cx = int(bx) + int(bw) // 2; cy = int(by) + int(bh) // 2
    # 가로 타원 — 원이 아닌 넓은 타원으로 infinity 모양 구현
    oval_w = int(bw * 0.46); oval_h = int(bh * 0.72)
    overlap = int(oval_w * 0.15)
    left_cx = cx - (oval_w - overlap) // 2; right_cx = cx + (oval_w - overlap) // 2

    from pptx.oxml.ns import qn as _qn
    def _set_alpha(shape, val='30000'):
        try:
            sf = shape.fill._fill.find(_qn('a:solidFill'))
            if sf is not None:
                sc = sf.find(_qn('a:srgbClr'))
                if sc is not None:
                    a_el = sc.makeelement(_qn('a:alpha'), {}); a_el.set('val', val); sc.append(a_el)
        except Exception: pass

    # 왼쪽 타원 (반투명)
    left_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_cx - oval_w // 2, cy - oval_h // 2, oval_w, oval_h)
    left_oval.fill.solid(); left_oval.fill.fore_color.rgb = COLORS["SEM_BLUE_BG"]
    left_oval.line.color.rgb = COLORS["PRIMARY"]; left_oval.line.width = Pt(2.5)
    _set_alpha(left_oval)

    # 오른쪽 타원 (반투명)
    right_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_cx - oval_w // 2, cy - oval_h // 2, oval_w, oval_h)
    right_oval.fill.solid(); right_oval.fill.fore_color.rgb = COLORS["SEM_GREEN_BG"]
    right_oval.line.color.rgb = COLORS["SEM_GREEN"]; right_oval.line.width = Pt(2.5)
    _set_alpha(right_oval)

    # 루프 라벨 (각 타원 바로 위에 배치 — 내부 항목과 겹침 방지)
    label_w = Inches(2.0); label_h = Inches(0.38)
    if left_label:
        ll_y = cy - oval_h // 2 - int(label_h) - Inches(0.03)
        ll_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_cx - int(label_w) // 2, ll_y, int(label_w), int(label_h))
        ll_shp.fill.solid(); ll_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; ll_shp.line.color.rgb = COLORS["PRIMARY"]
        tf = ll_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = left_label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(12); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER
    if right_label:
        rl_y = cy - oval_h // 2 - int(label_h) - Inches(0.03)
        rl_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_cx - int(label_w) // 2, rl_y, int(label_w), int(label_h))
        rl_shp.fill.solid(); rl_shp.fill.fore_color.rgb = COLORS["SEM_GREEN"]; rl_shp.line.color.rgb = COLORS["SEM_GREEN"]
        tf = rl_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = right_label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(12); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    # 중앙 교차 라벨
    if center_label:
        c_w = Inches(1.3); c_h = Inches(0.55)
        c_shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - int(c_w) // 2, cy - int(c_h) // 2, int(c_w), int(c_h))
        c_shp.fill.solid(); c_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; c_shp.line.color.rgb = COLORS["BG_WHITE"]; c_shp.line.width = Pt(2.5)
        tf = c_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = center_label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(12); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    # 항목을 각 타원의 바깥 호를 따라 배치 (angle-based)
    item_w = Inches(1.5); item_h = Inches(0.45)
    semi_a = oval_w // 2 - int(item_w) // 2 - Inches(0.12)
    semi_b = oval_h // 2 - int(item_h) // 2 - Inches(0.12)

    # 왼쪽 항목: 바깥쪽 호 (top → upper-left → lower-left → bottom)
    left_angles = [-math.pi / 2, -5 * math.pi / 6, 5 * math.pi / 6, math.pi / 2]
    n_left = min(len(left_items), 4)
    for i in range(n_left):
        angle = left_angles[i]
        ix = left_cx + int(semi_a * math.cos(angle)) - int(item_w) // 2
        iy = cy + int(semi_b * math.sin(angle)) - int(item_h) // 2
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ix, iy, int(item_w), int(item_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = COLORS["BG_WHITE"]
        shp.line.color.rgb = COLORS["PRIMARY"]; shp.line.width = Pt(1.5)
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        label = left_items[i].get('label', '') if isinstance(left_items[i], dict) else str(left_items[i])
        p = tf.paragraphs[0]; p.text = label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(11); p.font.color.rgb = COLORS["PRIMARY"]; p.alignment = PP_ALIGN.CENTER

    # 오른쪽 항목: 바깥쪽 호 (top → upper-right → lower-right → bottom)
    right_angles = [-math.pi / 2, -math.pi / 6, math.pi / 6, math.pi / 2]
    n_right = min(len(right_items), 4)
    for i in range(n_right):
        angle = right_angles[i]
        ix = right_cx + int(semi_a * math.cos(angle)) - int(item_w) // 2
        iy = cy + int(semi_b * math.sin(angle)) - int(item_h) // 2
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ix, iy, int(item_w), int(item_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = COLORS["BG_WHITE"]
        shp.line.color.rgb = COLORS["SEM_GREEN"]; shp.line.width = Pt(1.5)
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        label = right_items[i].get('label', '') if isinstance(right_items[i], dict) else str(right_items[i])
        p = tf.paragraphs[0]; p.text = label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(11); p.font.color.rgb = COLORS["SEM_GREEN"]; p.alignment = PP_ALIGN.CENTER

    # 흐름 방향 화살표: 각 아이템 사이 + 교차 지점
    def _add_flow_dot(fx, fy, rot, color):
        """작은 원에 ➤ 방향 표시"""
        ds = Inches(0.28)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, fx - int(ds) // 2, fy - int(ds) // 2, int(ds), int(ds))
        d.fill.solid(); d.fill.fore_color.rgb = color; d.line.color.rgb = color
        d.rotation = rot
        tf_d = d.text_frame; tf_d.clear(); tf_d.vertical_anchor = MSO_ANCHOR.MIDDLE
        pd = tf_d.paragraphs[0]; pd.text = "\u27A4"; pd.font.size = Pt(9); pd.font.color.rgb = COLORS["BG_WHITE"]; pd.alignment = PP_ALIGN.CENTER

    # 왼쪽 루프 흐름 (반시계: top→upper-left→lower-left→bottom)
    for j in range(n_left - 1):
        a1 = left_angles[j]; a2 = left_angles[j + 1]
        mid_a = (a1 + a2) / 2
        mx = left_cx + int(semi_a * 0.75 * math.cos(mid_a))
        my = cy + int(semi_b * 0.75 * math.sin(mid_a))
        _add_flow_dot(mx, my, math.degrees(mid_a) + 90, COLORS["PRIMARY"])

    # 오른쪽 루프 흐름 (시계: top→upper-right→lower-right→bottom)
    for j in range(n_right - 1):
        a1 = right_angles[j]; a2 = right_angles[j + 1]
        mid_a = (a1 + a2) / 2
        mx = right_cx + int(semi_a * 0.75 * math.cos(mid_a))
        my = cy + int(semi_b * 0.75 * math.sin(mid_a))
        _add_flow_dot(mx, my, math.degrees(mid_a) + 90, COLORS["SEM_GREEN"])

    # 교차 화살표 (큰 쉐브론)
    arr_w = Inches(0.55); arr_h = Inches(0.32)
    a1_shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + Inches(0.08), cy + Inches(0.35), int(arr_w), int(arr_h))
    a1_shp.fill.solid(); a1_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; a1_shp.line.fill.background()
    a2_shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, cx - Inches(0.08) - int(arr_w), cy - Inches(0.35) - int(arr_h), int(arr_w), int(arr_h))
    a2_shp.rotation = 180; a2_shp.fill.solid(); a2_shp.fill.fore_color.rgb = COLORS["SEM_GREEN"]; a2_shp.line.fill.background()


# 36. Mind Map (마인드맵)
def render_mind_map(slide, data):
    import math
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    center = content.get('center', {}); branches = content.get('branches', [])
    if not center: return
    n = len(branches)
    if not n: return
    default_colors = ['blue', 'green', 'orange', 'red', 'primary', 'gray']

    # ── 좌측 55% = 방사형 마인드맵 ──
    left_w = int(int(bw) * 0.52)
    map_cx = int(bx) + left_w // 2
    map_cy = int(by) + int(bh) // 2

    # 중앙 원 (OVAL)
    center_r = int(min(int(bh), left_w) * 0.16)
    center_d = center_r * 2
    c_shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, map_cx - center_r, map_cy - center_r, center_d, center_d)
    c_shp.fill.solid(); c_shp.fill.fore_color.rgb = COLORS["PRIMARY"]; c_shp.line.color.rgb = COLORS["PRIMARY"]
    tf = c_shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    p = tf.paragraphs[0]; p.text = center.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(11); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    # 브랜치 노드 (라운드 사각형 — 라벨 + sub_branches 표시)
    node_w = Inches(1.4); node_h = Inches(0.7)
    nhw = int(node_w) // 2; nhh = int(node_h) // 2
    half_lw = left_w // 2; half_bh = int(bh) // 2

    # 균일 선 길이 계산 — 좌측 영역 경계 제약
    max_ll = int(Inches(20))
    for i in range(n):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        ca = abs(math.cos(angle)); sa = abs(math.sin(angle))
        if ca < 0.01: dr = nhh
        elif sa < 0.01: dr = nhw
        else: dr = min(int(nhw / ca), int(nhh / sa))
        if ca > 0.01:
            max_ll = min(max_ll, int((half_lw - nhw - Inches(0.08)) / ca) - dr - center_r)
        if sa > 0.01:
            max_ll = min(max_ll, int((half_bh - nhh - Inches(0.08)) / sa) - dr - center_r)
    line_len = max(Inches(0.2), max_ll)

    for i, br in enumerate(branches):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        cos_a = math.cos(angle); sin_a = math.sin(angle)
        aca = abs(cos_a); asa = abs(sin_a)
        if aca < 0.01: dr = nhh
        elif asa < 0.01: dr = nhw
        else: dr = min(int(nhw / aca), int(nhh / asa))

        color_key = br.get('color', default_colors[i % len(default_colors)])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        # 커넥터: 원 가장자리 → 노드 면
        ex = map_cx + int(center_r * cos_a); ey = map_cy + int(center_r * sin_a)
        face_x = map_cx + int((center_r + line_len) * cos_a)
        face_y = map_cy + int((center_r + line_len) * sin_a)
        conn = slide.shapes.add_connector(1, ex, ey, face_x, face_y)
        conn.line.color.rgb = line_c; conn.line.width = Pt(2.0)

        # 노드 박스 중심
        ncx = map_cx + int((center_r + line_len + dr) * cos_a)
        ncy = map_cy + int((center_r + line_len + dr) * sin_a)
        ns = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ncx - nhw, ncy - nhh, int(node_w), int(node_h))
        ns.fill.solid(); ns.fill.fore_color.rgb = fill_c; ns.line.color.rgb = line_c; ns.line.width = Pt(1.5)
        ntf = ns.text_frame; ntf.clear(); ntf.word_wrap = True; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ntf.margin_left = Inches(0.05); ntf.margin_right = Inches(0.05); ntf.margin_top = Inches(0.03)
        p = ntf.paragraphs[0]; p.text = br.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(9); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(1)
        for sub in br.get('sub_branches', [])[:2]:
            p2 = ntf.add_paragraph(); p2.text = f"· {sub}"; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(7); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(0)

    # ── 우측 48% = 설명 카드 ──
    right_x = int(bx) + left_w + Inches(0.2)
    right_w = int(bw) - left_w - Inches(0.2)
    card_gap = Inches(0.1)
    card_h = (int(bh) - int(card_gap) * max(n - 1, 1)) // max(n, 1)
    bar_w = Inches(0.06)

    for i, br in enumerate(branches):
        color_key = br.get('color', default_colors[i % len(default_colors)])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])
        cy = int(by) + int((card_h + int(card_gap)) * i)

        # 색상 악센트 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, cy, int(bar_w), card_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = line_c; bar.line.fill.background()

        # 카드
        card_x = right_x + int(bar_w); cw = right_w - int(bar_w)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, cy, cw, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = fill_c; card.line.color.rgb = COLORS["BORDER"]; card.line.width = Pt(0.75)
        tf = card.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.08); tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]; p.text = br.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(10); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(2)
        desc = br.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(8); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.LEFT; p2.space_after = Pt(1)
        else:
            for sub in br.get('sub_branches', [])[:3]:
                p2 = tf.add_paragraph(); p2.text = f"• {sub}"; p2.font.name = FONTS["BODY_TEXT"]
                p2.font.size = Pt(8); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.LEFT; p2.space_after = Pt(1)


# 38. Checklist 2-Column (2열 체크리스트 그리드)
def render_checklist_2col(slide, data):
    """2열 체크리스트 그리드 — 원본 p6 구조 그대로 (plain TextBox 요약, 고정 0.30 서브항목 행, 배경 없음)

    data.data.data:
        summary: "1/10 Passed  9 Warning"  # 선택, 상단 plain TextBox 18pt bold
        items: [
            {
                "title": "WBS 1.1 태스크 헤더 (status emoji 포함 가능)",
                "status": "done" | "in_progress" | "todo",
                "subitems": [
                    {"text": "하위 항목 텍스트", "badge": "CRITICAL" | "HIGH" | "MEDIUM" | ""}
                ]
            }, ...
        ]
    최대 6개 items (3행×2열), 각 item당 subitems 최대 3개 권장
    col_gap=0.25in, subitem row 고정 0.30in, 배경 Rectangle 없음
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content_data = wrapper.get('data', {})
    items = content_data.get('items', [])
    if not items: return

    # 요약 텍스트 — plain TextBox (18pt bold, 배경 없음) — 원본 p6 구조
    summary_text = content_data.get('summary', '')
    if summary_text:
        sum_h = Inches(0.40)
        stb = slide.shapes.add_textbox(bx, by, bw, sum_h)
        stf = stb.text_frame; stf.word_wrap = False; stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = stf.paragraphs[0]; sp.text = summary_text
        sp.font.name = FONTS["BODY_TEXT"]; sp.font.size = Pt(18); sp.font.bold = True
        sp.font.color.rgb = COLORS["DARK_GRAY"]; sp.alignment = PP_ALIGN.LEFT
        by = by + sum_h + Inches(0.05)
        bh = bh - sum_h - Inches(0.05)

    # 진행률 바 (원본 p6: bg=F8F9FA/DCDCDC, green=047857, orange=C2410C, h=0.28)
    # summary 텍스트에서 비율 파싱 (예: "1/10 Passed 9 Warning" → done=1, total=10)
    import re as _re
    _m = _re.search(r'(\d+)/(\d+)', summary_text) if summary_text else None
    if _m:
        n_done = int(_m.group(1)); n_total = int(_m.group(2))
    else:
        n_done = sum(1 for it in items if it.get('status') == 'done'); n_total = max(len(items), 1)
    n_warn = n_total - n_done
    progress_h = Inches(0.28)
    pb_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, progress_h)
    pb_bg.fill.solid(); pb_bg.fill.fore_color.rgb = COLORS["BG_BOX"]
    pb_bg.line.color.rgb = COLORS["BORDER"]; pb_bg.line.width = Pt(1.0)
    green_w = int(bw * n_done / n_total)
    if green_w > 0:
        pb_g = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, green_w, progress_h)
        pb_g.fill.solid(); pb_g.fill.fore_color.rgb = COLORS["SEM_GREEN"]; pb_g.line.fill.background()
    if n_warn > 0:
        pb_w = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx + green_w, by, bw - green_w, progress_h)
        pb_w.fill.solid(); pb_w.fill.fore_color.rgb = COLORS["SEM_ORANGE"]; pb_w.line.fill.background()
    by = by + progress_h + Inches(0.27)
    bh = bh - progress_h - Inches(0.27)

    # 원본: col_gap=0.250in, col_w=6.041in
    col_gap = Inches(0.25); col_w = int((bw - col_gap) / 2)
    n_items = len(items); n_rows = (n_items + 1) // 2

    # subitem 행 고정 0.30in, header 0.35in, gap 0.04in
    hdr_h = Inches(0.35); sub_row_h = Inches(0.30); sub_gap = Inches(0.04)
    max_subs = max((len(item.get('subitems', [])) for item in items), default=1)
    row_h_natural = int(hdr_h + sub_gap + max_subs * sub_row_h)
    row_gap = Inches(0.10)
    total_nat = row_h_natural * n_rows + int(row_gap) * (n_rows - 1)
    row_h = int(row_h_natural * bh / total_nat) if total_nat > bh else row_h_natural

    # 아이콘: 상태별, 배지: 종류별
    status_icons  = {"done": "✓", "in_progress": "⚠", "todo": "○"}
    icon_colors   = {"done": COLORS["SEM_GREEN"], "in_progress": COLORS["SEM_ORANGE"], "todo": COLORS["DARK_GRAY"]}
    badge_colors  = {
        "CRITICAL": (COLORS["SEM_RED"],    COLORS["SEM_RED_BG"]),
        "HIGH":     (COLORS["SEM_ORANGE"], COLORS["SEM_ORANGE_BG"]),
        "MEDIUM":   (COLORS["SEM_BLUE"],   COLORS["SEM_BLUE_BG"]),
    }

    for idx, item in enumerate(items):
        col = idx % 2; row = idx // 2
        cx = bx + col * (col_w + col_gap)
        cy = by + row * (row_h + row_gap)

        # 헤더 — 원본 p6: 모든 항목 동일 SEM_BLUE(1E3A8A) fill, 0pt border, 흰색 텍스트
        status = item.get('status', 'todo')
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, col_w, hdr_h)
        hdr.fill.solid(); hdr.fill.fore_color.rgb = COLORS["SEM_BLUE"]
        hdr.line.color.rgb = COLORS["SEM_BLUE"]; hdr.line.width = Pt(0)
        tf = hdr.text_frame; tf.clear(); tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = item.get('title', '')
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(11)
        p.font.color.rgb = COLORS["BG_WHITE"]

        # 서브 항목들 (배경 Rectangle 없음 — 원본 p6 구조)
        subitems = item.get('subitems', [])[:4]
        sub_y = cy + hdr_h + sub_gap
        icon = status_icons.get(status, '○')
        icon_c = icon_colors.get(status, COLORS["DARK_GRAY"])
        bdg_w = Inches(0.650); bdg_h_val = Inches(0.220)
        txt_fixed_w = Inches(4.891)

        for si, sub in enumerate(subitems):
            sy = sub_y + si * sub_row_h
            badge_key = sub.get('badge', '')
            badge_pair = badge_colors.get(badge_key, None)

            # 상태 아이콘 TextBox (0.300×0.300)
            itb = slide.shapes.add_textbox(cx + Inches(0.080), sy, Inches(0.300), sub_row_h)
            ift = itb.text_frame; ift.word_wrap = False; ift.margin_left = Inches(0); ift.margin_top = Inches(0)
            ift.vertical_anchor = MSO_ANCHOR.MIDDLE
            ip = ift.paragraphs[0]; ip.text = icon
            ip.font.name = FONTS["BODY_TEXT"]; ip.font.size = Pt(10); ip.font.color.rgb = icon_c; ip.alignment = PP_ALIGN.CENTER

            # 텍스트 TextBox (4.891×0.300 @ cx+0.420)
            tb = slide.shapes.add_textbox(cx + Inches(0.420), sy, txt_fixed_w, sub_row_h)
            tf2 = tb.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0); tf2.margin_top = Inches(0)
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]; p2.text = sub.get('text', '')
            p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(10); p2.font.color.rgb = COLORS["DARK_GRAY"]

            # 배지 TextBox (0.650×0.220 @ cx+col_w-0.730, y+0.040)
            if badge_pair:
                bc, bb = badge_pair
                bdg_x = cx + col_w - bdg_w - Inches(0.080)
                bdg_y = sy + Inches(0.040)
                bdg = slide.shapes.add_textbox(bdg_x, bdg_y, bdg_w, bdg_h_val)
                bdg.fill.solid(); bdg.fill.fore_color.rgb = bb
                tf3 = bdg.text_frame; tf3.clear(); tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
                p3 = tf3.paragraphs[0]; p3.text = badge_key
                p3.font.name = FONTS["BODY_TEXT"]; p3.font.size = Pt(8); p3.font.bold = True
                p3.font.color.rgb = bc; p3.alignment = PP_ALIGN.CENTER


# 39. Kanban Board (칸반 보드 — 열당 N개 개별 카드)
def render_kanban_board(slide, data):
    """Kanban 보드 — 원본 p7 구조 그대로
    컬럼 전체높이 border + 헤더 + (RoundedRect+좌측accent bar+TextBox) 카드 구조

    data.data.data.columns: [
        {
            "title": "To Do (3)",
            "color": "gray" | "blue" | "green" | "orange" | "red",
            "cards": [
                {"title": "태스크 제목\n날짜, 담당, 기간", "badge": "Critical" | "내일" | "완료" | ""}
            ]
        }
    ]
    카드 고정 높이 0.70in, 최대 표시 카드 = 가용 높이에서 자동 계산
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    columns = wrapper.get('data', {}).get('columns', [])
    if not columns: return

    col_palette = {
        "navy":   (COLORS["SEM_BLUE"],     COLORS["SEM_BLUE_BG"]),   # 1E3A8A (To Do)
        "blue":   (COLORS["PRIMARY"],      COLORS["SEM_BLUE_BG"]),   # 0043DA
        "green":  (COLORS["SEM_GREEN"],    COLORS["SEM_GREEN_BG"]),
        "orange": (COLORS["SEM_ORANGE"],   COLORS["SEM_ORANGE_BG"]),
        "red":    (COLORS["SEM_RED"],      COLORS["SEM_RED_BG"]),
        "gray":   (COLORS["DARK_GRAY"],    COLORS["BG_BOX"]),
    }
    n_cols = len(columns); col_gap = Inches(0.15)
    col_w = int((bw - col_gap * (n_cols - 1)) / n_cols)
    hdr_h = Inches(0.5); card_h = Inches(0.70); card_gap = Inches(0.10)
    card_margin = Inches(0.080); accent_bar_w = Inches(0.060)
    card_inner_w = col_w - card_margin * 2
    txt_x_off = accent_bar_w + Inches(0.100)
    txt_w = card_inner_w - txt_x_off - Inches(0.10)  # 우측 margin 0.10 (원본 p7)
    txt_h = Inches(0.580)
    bdg_w = Inches(0.650); bdg_h = Inches(0.200)
    available_card_h = bh - hdr_h - Inches(0.10)     # header 아래 gap 0.10 (원본 p7)
    max_cards = max(int(available_card_h / (card_h + card_gap)), 1)

    for ci, col in enumerate(columns):
        cx = bx + ci * (col_w + col_gap)
        accent, _ = col_palette.get(col.get('color', 'navy'), col_palette['navy'])

        # 컨테이너 (원본 p7: fill=F8F9FA, line=DCDCDC/1pt)
        container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, by, col_w, bh)
        container.fill.solid(); container.fill.fore_color.rgb = COLORS["BG_BOX"]
        container.line.color.rgb = COLORS["BORDER"]; container.line.width = Pt(1.0)

        # 열 헤더 (원본 p7: fill=accent solid, line=accent/0pt, 흰색 텍스트)
        col_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, by, col_w, hdr_h)
        col_hdr.fill.solid(); col_hdr.fill.fore_color.rgb = accent
        col_hdr.line.color.rgb = accent; col_hdr.line.width = Pt(0)
        tf = col_hdr.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = col.get('title', '')
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(13)
        p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

        cards = col.get('cards', [])[:max_cards]
        card_y = by + hdr_h + Inches(0.10)

        for card in cards:
            card_x = cx + card_margin

            # 카드 (원본 p7: fill=FFFFFF, line=DCDCDC/1pt)
            crr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_inner_w, card_h)
            crr.fill.solid(); crr.fill.fore_color.rgb = COLORS["BG_WHITE"]
            crr.line.color.rgb = COLORS["BORDER"]; crr.line.width = Pt(1.0)

            # 좌측 accent bar
            abar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x, card_y, accent_bar_w, card_h)
            abar.fill.solid(); abar.fill.fore_color.rgb = accent; abar.line.fill.background()

            # 카드 텍스트
            tb = slide.shapes.add_textbox(card_x + txt_x_off, card_y + Inches(0.060), txt_w, txt_h)
            tf2 = tb.text_frame; tf2.word_wrap = True; tf2.vertical_anchor = MSO_ANCHOR.TOP
            tf2.margin_left = Inches(0); tf2.margin_top = Inches(0)
            for li, line_text in enumerate(card.get('title', '').split('\n')):
                p2 = tf2.paragraphs[0] if li == 0 else tf2.add_paragraph()
                p2.text = line_text; p2.font.name = FONTS["BODY_TEXT"]
                p2.font.size = Pt(10) if li == 0 else Pt(9)
                p2.font.color.rgb = COLORS["DARK_GRAY"]

            # 배지 (원본 p7: "Critical"→FEF2F2/B91C1C, 기타→F8F9FA/505050)
            badge_text = card.get('badge', '')
            if badge_text:
                is_crit = 'critical' in badge_text.lower()
                bdg_fill = COLORS["SEM_RED_BG"] if is_crit else COLORS["BG_BOX"]
                bdg_line = COLORS["SEM_RED"]    if is_crit else COLORS["GRAY"]
                bdg_x = card_x + card_inner_w - bdg_w - Inches(0.080)
                bdg_y_pos = card_y + card_h - bdg_h - Inches(0.080)
                bdg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bdg_x, bdg_y_pos, bdg_w, bdg_h)
                bdg.fill.solid(); bdg.fill.fore_color.rgb = bdg_fill
                bdg.line.color.rgb = bdg_line; bdg.line.width = Pt(1.0)
                tf3 = bdg.text_frame; tf3.clear(); tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf3.margin_left = Inches(0.02); tf3.margin_right = Inches(0.02)  # 원본 p7: ml=mr=0.02
                tf3.margin_top = Inches(0.05); tf3.margin_bottom = Inches(0.05)
                p3 = tf3.paragraphs[0]; p3.text = badge_text
                p3.font.name = FONTS["BODY_TEXT"]; p3.font.size = Pt(8); p3.font.bold = True
                p3.font.color.rgb = bdg_line; p3.alignment = PP_ALIGN.CENTER

            card_y += card_h + card_gap


# 40. Executive Summary (전체 폭 레이블+본문 섹션)
def render_exec_summary(slide, data):
    """Executive Summary — 원본 p9 구조 그대로

    섹션 타입 자동 감지 (body 텍스트 기준):
    - '■' 포함 → stacked, ■ 줄별 separate TextBox (10pt)
    - '➤' 포함 → stacked + left accent bar + ONE TextBox (11pt)
    - 나머지   → inline, label 좌 / body 우 side-by-side (11pt)

    data.data.data.sections: [
        {"label": "상황",        "body": "...",          "color": "gray"},
        {"label": "핵심 발견사항", "body": "■ L1\\n■ L2",  "color": "blue"},
        {"label": "권고사항",     "body": "➤ L1\\n➤ L2",  "color": "orange"},
    ]
    레이블 배지 너비: 2자=0.90in, 이후 글자당 +0.10in
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    sections = wrapper.get('data', {}).get('sections', [])
    if not sections: return

    # 원본 p9 색상 매핑 — 섹션 배경 border는 항상 DCDCDC
    color_map = {
        "gray":   (COLORS["GRAY"],     COLORS["BG_BOX"],      COLORS["BORDER"]),  # 505050, F8F9FA, DCDCDC
        "blue":   (COLORS["PRIMARY"],  COLORS["BG_WHITE"],    COLORS["BORDER"]),  # 0043DA, FFFFFF, DCDCDC
        "navy":   (COLORS["SEM_BLUE"], COLORS["SEM_BLUE_BG"], COLORS["BORDER"]),  # 1E3A8A, EFF6FF, DCDCDC
        "green":  (COLORS["SEM_GREEN"],  COLORS["SEM_GREEN_BG"],  COLORS["BORDER"]),
        "orange": (COLORS["SEM_ORANGE"], COLORS["SEM_ORANGE_BG"], COLORS["BORDER"]),
        "red":    (COLORS["SEM_RED"],    COLORS["SEM_RED_BG"],    COLORS["BORDER"]),
    }

    # 원본 p9 측정값 기반 상수
    lbl_h = Inches(0.280); content_gap = Inches(0.080)
    bottom_pad = Inches(0.120); bullet_plain_h = Inches(0.500); bullet_arrow_h = Inches(0.290)
    inline_h = Inches(0.946); sec_gap = Inches(0.150); accent_bar_w = Inches(0.080)
    # 타입별 badge top padding (원본 p9 실측):
    #   inline/bullet_plain: badge y-offset = 0.120" (shape[5] y=2.82-2.7=0.12, shape[8] y=3.866-3.746=0.12)
    #   bullet_arrow: badge y-offset = 0.100" (shape[14] y=5.795-5.695=0.10) → 좌측 accent bar 때문에 다름
    lbl_top_pad       = Inches(0.120)  # inline / bullet_plain
    lbl_top_pad_arrow = Inches(0.100)  # bullet_arrow 전용

    def _layout_type(body_text):
        if '■' in body_text: return 'bullet_plain'
        if '➤' in body_text: return 'bullet_arrow'
        return 'inline'

    def _ideal_h(body_text):
        lt = _layout_type(body_text)
        if lt == 'inline': return int(inline_h)
        if lt == 'bullet_plain':
            # ■ 로 시작하는 줄만 쌍(pair)으로 카운트
            n = max(sum(1 for l in body_text.split('\n') if l.strip().startswith('■')), 1)
            return int(lbl_top_pad + lbl_h + content_gap + n * bullet_plain_h + bottom_pad)
        else:  # bullet_arrow — lbl_top_pad_arrow(0.100) 사용
            n = max(len([l for l in body_text.split('\n') if l.strip()]), 1)
            return int(lbl_top_pad_arrow + lbl_h + content_gap + n * bullet_arrow_h + bottom_pad)

    # 이상적 높이 → 비례 스케일
    n = len(sections)
    bodies = [s.get('body', '') for s in sections]
    ideal_heights = [_ideal_h(b) for b in bodies]
    total_ideal = sum(ideal_heights) + sec_gap * (n - 1)
    scale = bh / total_ideal if total_ideal > bh else 1.0
    sec_heights = [int(h * scale) for h in ideal_heights]

    for i, sec in enumerate(sections):
        sy = by + sum(sec_heights[:i]) + sec_gap * i
        sec_h = sec_heights[i]
        body_text = sec.get('body', '')
        lt = _layout_type(body_text)
        accent, bg, txt_c = color_map.get(sec.get('color', 'gray'), color_map['gray'])

        # 배경 전체 바 (원본 p9: border는 항상 DCDCDC)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, sy, bw, sec_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = bg
        bar.line.color.rgb = COLORS["BORDER"]; bar.line.width = Pt(1.0)

        # ➤ 타입: 좌측 accent bar (원본 p9: fill=PRIMARY=0043DA, 섹션 accent와 무관)
        if lt == 'bullet_arrow':
            abar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, sy, accent_bar_w, sec_h)
            abar.fill.solid(); abar.fill.fore_color.rgb = COLORS["PRIMARY"]; abar.line.fill.background()
            lbl_x = bx + Inches(0.250)
            cur_lbl_top_pad = lbl_top_pad_arrow  # 원본: 0.100"
        else:
            lbl_x = bx + Inches(0.150)
            cur_lbl_top_pad = lbl_top_pad         # 원본: 0.120"

        # 레이블 배지: 너비=label 길이 기반, y=sy+cur_lbl_top_pad
        label_text = sec.get('label', '')
        lbl_w = Inches(0.900 + max(0, len(label_text) - 2) * 0.100)
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lbl_x, sy + cur_lbl_top_pad, lbl_w, lbl_h)
        lbl.fill.solid(); lbl.fill.fore_color.rgb = accent
        lbl.line.color.rgb = accent; lbl.line.width = Pt(0)
        tf = lbl.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)  # 원본 p9: ml=mr=0.02
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)  # 원본 p9: mt=mb=0.05
        p = tf.paragraphs[0]; p.text = label_text
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(10)  # 원본: 10pt
        p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

        # 본문 시작 y (타입별 cur_lbl_top_pad 적용)
        body_y = sy + cur_lbl_top_pad + lbl_h + content_gap

        if lt == 'inline':
            # label 우측 body — lbl_x + lbl_w + gap 이후 시작해서 배지와 겹침 방지
            lbl_right = lbl_x + lbl_w + Inches(0.150)
            body_x = lbl_right
            body_w_val = bx + bw - lbl_right - Inches(0.150)
            btb = slide.shapes.add_textbox(body_x, sy + Inches(0.080), body_w_val, sec_h - Inches(0.100))
            tf2 = btb.text_frame; tf2.word_wrap = True; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf2.margin_left = Inches(0.05); tf2.margin_right = Inches(0.05)  # 원본 p9: ml=mr=0.05
            tf2.margin_top = Inches(0.05); tf2.margin_bottom = Inches(0.05)  # 원본 p9: mt=mb=0.05
            p2 = tf2.paragraphs[0]; p2.text = body_text
            p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(11); p2.font.bold = False
            p2.font.color.rgb = COLORS["DARK_GRAY"]

        elif lt == 'bullet_plain':
            # ■ heading+body 쌍별 separate TextBox (원본 p9: x=0.750, w=11.833)
            # 데이터 형식: "■  heading1\n    body1\n■  heading2\n    body2"
            body_x = bx + Inches(0.250)  # 원본: x=0.75, bx=0.5 → offset=0.25
            body_w_val = bw - Inches(0.250) - Inches(0.250)  # 원본: w=11.833
            raw_lines = body_text.split('\n')
            # (heading, body) 쌍 파싱
            pairs = []
            i2 = 0
            while i2 < len(raw_lines):
                rl = raw_lines[i2]
                if rl.strip().startswith('■'):
                    heading = rl
                    body_ln = raw_lines[i2+1] if i2+1 < len(raw_lines) and not raw_lines[i2+1].strip().startswith('■') else ''
                    pairs.append((heading, body_ln))
                    i2 += 2 if body_ln else 1
                else:
                    i2 += 1
            for li, (hdr_line, bdy_line) in enumerate(pairs):
                ly = body_y + li * bullet_plain_h
                ltb = slide.shapes.add_textbox(body_x, ly, body_w_val, bullet_plain_h)
                tf3 = ltb.text_frame; tf3.word_wrap = True; tf3.vertical_anchor = MSO_ANCHOR.TOP
                tf3.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 크기 원본과 동일, 넘칠 경우 폰트 축소
                tf3.margin_left = Inches(0.05); tf3.margin_right = Inches(0.05)  # 원본 p9: ml=mr=0.05
                tf3.margin_top = Inches(0.02); tf3.margin_bottom = Inches(0.02)  # 원본 p9: mt=mb=0.02
                p3 = tf3.paragraphs[0]; p3.text = hdr_line
                p3.font.name = FONTS["BODY_TEXT"]; p3.font.size = Pt(11); p3.font.bold = True  # 원본: 11pt bold
                p3.font.color.rgb = COLORS["DARK_GRAY"]  # 원본: 212121
                p3.space_after = Pt(4)  # heading과 body 간격
                if bdy_line.strip():
                    p3b = tf3.add_paragraph(); p3b.text = bdy_line
                    p3b.font.name = FONTS["BODY_TEXT"]; p3b.font.size = Pt(9); p3b.font.bold = False  # 원본: 9pt
                    p3b.font.color.rgb = COLORS["GRAY"]  # 원본: 505050

        else:  # bullet_arrow
            # ➤ ONE TextBox (원본 p9: x=0.800, w=11.733, h=0.945)
            body_x = bx + Inches(0.300)  # 원본: x=0.80
            body_w_val = bw - Inches(0.300) - Inches(0.300)  # 원본: w=11.733
            # cur_lbl_top_pad=0.100 으로 body_h_val = sec_h - 0.100 - 0.280 - 0.080 - 0.100 = sec_h - 0.560
            # scale=1.0 이면 sec_h=1.505" → body_h_val=0.945" = 3×0.315" 정확히 일치
            body_h_val = sec_h - cur_lbl_top_pad - lbl_h - content_gap - bottom_pad
            lines = [l for l in body_text.split('\n') if l.strip()]
            # 하한선 보장: 공간이 부족하면 들어갈 수 있는 줄 수로만 축소
            max_lines_fit = max(1, int(body_h_val // bullet_arrow_h))
            if len(lines) > max_lines_fit:
                lines = lines[:max_lines_fit]
            btb = slide.shapes.add_textbox(body_x, body_y, body_w_val, body_h_val)
            tf4 = btb.text_frame; tf4.word_wrap = True; tf4.vertical_anchor = MSO_ANCHOR.TOP
            tf4.margin_left = Inches(0.05); tf4.margin_right = Inches(0.05)  # 원본 p9: ml=mr=0.05
            tf4.margin_top = Inches(0.02); tf4.margin_bottom = Inches(0.05)  # 원본 p9: mt=0.02, mb=0.05
            for li, line in enumerate(lines):
                p4 = tf4.paragraphs[0] if li == 0 else tf4.add_paragraph()
                p4.text = line
                p4.font.name = FONTS["BODY_TEXT"]; p4.font.size = Pt(10); p4.font.bold = False  # 원본: 10pt
                p4.font.color.rgb = COLORS["SEM_BLUE_TEXT"]; p4.space_after = Pt(4)  # 원본: 1E40AF


# 41. Risk Table (원본 p8 — 리스크 & 전제 조건)
def render_risk_table(slide, data):
    """위험/점검 항목 테이블 — 원본 p8 구조 그대로

    상단 요약 바 (SEM_BLUE) + 컬럼 헤더 + 구분선 + 교대 색 행
    행: 상태 circle (C2410C=high, B91C1C=critical) | 항목 | 설명 | 담당자

    data.data.data:
      summary: "⬤ 4 Yellow   |   ⬤ 1 Red"
      columns: ["상태", "항목", "설명", "담당자"]
      col_widths: [0.6, 3.083, 4.933, 3.717]  # optional, inches
      rows:
        - level: "high" | "critical"  → circle color
          item: "항목명"
          desc: "설명 텍스트"
          owner: "담당자"
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})
    rows_data = content.get('rows', [])
    if not rows_data: return

    # 원본 p8 측정값 기반 상수
    sum_bar_h  = Inches(0.550)  # 상단 요약 바 h
    col_hdr_y_off = Inches(0.150)  # 요약 바 아래 gap → col header y
    col_hdr_h  = Inches(0.400)  # col header row h
    divider_h  = Inches(0.014)  # 구분선 h
    row_h      = Inches(0.580)  # 데이터 행 h
    circle_sz  = Inches(0.319)  # 상태 circle w=h

    # 열 x 오프셋 / 너비 (원본 p8 측정값)
    col_widths_in = content.get('col_widths', [0.6, 3.083, 4.933, 3.717])
    col_xs = []
    cx = 0.0
    for cw in col_widths_in:
        col_xs.append(cx)
        cx += cw
    # Inch 변환
    col_xs_emu = [bx + Inches(x) for x in col_xs]
    col_ws_emu = [Inches(w) for w in col_widths_in]

    # 1. 상단 요약 바
    summary_text = content.get('summary', '')
    sum_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, sum_bar_h)
    sum_bar.fill.solid(); sum_bar.fill.fore_color.rgb = COLORS["SEM_BLUE"]
    sum_bar.line.color.rgb = COLORS["SEM_BLUE"]; sum_bar.line.width = Pt(0)
    if summary_text:
        tf_s = sum_bar.text_frame; tf_s.clear(); tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_s = tf_s.paragraphs[0]; p_s.text = summary_text
        p_s.font.name = FONTS["BODY_TEXT"]; p_s.font.size = Pt(13); p_s.font.bold = True
        p_s.font.color.rgb = COLORS["BG_WHITE"]; p_s.alignment = PP_ALIGN.CENTER

    # 2. 컬럼 헤더 TextBoxes
    col_hdr_y = by + sum_bar_h + col_hdr_y_off
    columns = content.get('columns', ['상태', '항목', '설명', '담당자'])
    for ci, col_name in enumerate(columns):
        if ci >= len(col_xs_emu): break
        tb = slide.shapes.add_textbox(col_xs_emu[ci], col_hdr_y, col_ws_emu[ci], col_hdr_h)
        tf = tb.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = col_name
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(12); p.font.bold = True
        p.font.color.rgb = COLORS["DARK_GRAY"]
        if ci == 0: p.alignment = PP_ALIGN.CENTER

    # 3. 구분선
    first_row_y = col_hdr_y + col_hdr_h
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, first_row_y - divider_h, bw, divider_h)
    div.fill.solid(); div.fill.fore_color.rgb = COLORS["BORDER"]; div.line.fill.background()

    # 4. 데이터 행들 (교대 FFFFFF / F8F9FA)
    row_fills = [COLORS["BG_WHITE"], COLORS["BG_BOX"]]
    for ri, row in enumerate(rows_data):
        ry = first_row_y + ri * row_h
        row_fill = row_fills[ri % 2]

        # 행 배경
        rbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, ry, bw, row_h)
        rbg.fill.solid(); rbg.fill.fore_color.rgb = row_fill; rbg.line.fill.background()

        # 상태 circle
        level = row.get('level', 'high')
        circle_c = COLORS["SEM_RED"] if level == 'critical' else (COLORS["SEM_ORANGE"] if level == 'orange' else COLORS["SEM_YELLOW"])
        c_x = col_xs_emu[0] + Inches(0.14)
        c_y = ry + (row_h - circle_sz) / 2
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, c_x, c_y, circle_sz, circle_sz)
        circ.fill.solid(); circ.fill.fore_color.rgb = circle_c
        circ.line.color.rgb = circle_c; circ.line.width = Pt(0)

        # 항목 / 설명 / 담당자 TextBoxes
        cell_vals = [row.get('item',''), row.get('desc',''), row.get('owner','')]
        for ci, val in enumerate(cell_vals):
            col_i = ci + 1  # 상태 열 스킵 → 항목(1), 설명(2), 담당자(3)
            if col_i >= len(col_xs_emu): break
            ctb = slide.shapes.add_textbox(col_xs_emu[col_i], ry, col_ws_emu[col_i], row_h)
            tf = ctb.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0); tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
            p = tf.paragraphs[0]; p.text = val
            p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(11)
            p.font.color.rgb = COLORS["DARK_GRAY"]


# ==========================================
# 5. 메인 라우터
# ==========================================
def render_speedometer_gauge(slide, data):
    """가로 바 게이지 — value(0-100), segments[], title
    세그먼트를 가로로 배치하고 value 위치에 마커 표시.
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    value = max(0, min(100, content.get('value', 0)))
    title = content.get('title', '')
    segments = content.get('segments', [])
    SLIDE_BOTTOM = Inches(7.4)

    seg_color_map = {
        'red':    RGBColor(220, 38,  38),
        'orange': RGBColor(234, 88,  12),
        'green':  RGBColor(22,  163, 74),
        'blue':   RGBColor(37,  99,  235),
        'gray':   RGBColor(156, 163, 175),
        'navy':   RGBColor(49,  46,  129),
    }
    n = max(len(segments), 1)

    # 레이아웃 상수 — 모두 by 기준 상대 위치
    bar_h      = Inches(0.55)
    bar_y      = int(by) + int(Inches(0.8))   # % 텍스트 공간 확보를 위해 아래로
    bar_x      = int(bx) + int(Inches(0.5))
    bar_w      = int(bw) - int(Inches(1.0))
    label_h    = Inches(0.35)
    marker_h   = Inches(0.30)
    value_h    = Inches(0.60)
    pct_h      = Inches(0.45)
    title_h    = Inches(0.40)

    # 슬라이드 경계 체크
    total_h = bar_h + label_h + Inches(0.15) + marker_h + Inches(0.1) + value_h + Inches(0.1) + pct_h + Inches(0.15) + title_h
    if bar_y + total_h > SLIDE_BOTTOM:
        bar_y = int(by) + int(Inches(0.1))

    # 1. 세그먼트 바
    seg_w = bar_w // n
    for i, seg in enumerate(segments):
        color_key = seg.get('color', 'gray')
        seg_rgb = seg_color_map.get(color_key, RGBColor(156, 163, 175))
        sx = bar_x + i * seg_w
        sw = seg_w if i < n - 1 else bar_w - i * seg_w  # 마지막은 나머지
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, bar_y, sw, int(bar_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = seg_rgb
        shp.line.color.rgb = COLORS["BG_WHITE"]; shp.line.width = Pt(1.5)

    # 2. 세그먼트 레이블 (바 아래)
    label_y = bar_y + int(bar_h) + int(Inches(0.05))
    for i, seg in enumerate(segments):
        color_key = seg.get('color', 'gray')
        seg_rgb = seg_color_map.get(color_key, RGBColor(156, 163, 175))
        sx = bar_x + i * seg_w
        sw = seg_w if i < n - 1 else bar_w - i * seg_w
        ltb = slide.shapes.add_textbox(sx, label_y, sw, int(label_h))
        tf = ltb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = seg.get('label', '')
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = seg_rgb; p.alignment = PP_ALIGN.CENTER

    # 3. value 마커 (바 위에 삼각형 + 선)
    marker_x = int(bar_x + bar_w * value / 100)
    marker_y = bar_y - int(Inches(0.05))
    # 세로선
    line = slide.shapes.add_connector(1, marker_x, bar_y - int(Inches(0.25)), marker_x, bar_y + int(bar_h) + int(Inches(0.25)))
    line.line.color.rgb = COLORS["DARK_GRAY"]; line.line.width = Pt(2.5)

    # 4. value % 텍스트 (마커 위)
    pct_y = bar_y - int(Inches(0.25)) - int(pct_h)
    pct_tb = slide.shapes.add_textbox(marker_x - int(Inches(0.6)), pct_y, int(Inches(1.2)), int(pct_h))
    tf = pct_tb.text_frame
    p = tf.paragraphs[0]; p.text = f"{value}%"
    p.font.name = FONTS["BODY_TITLE"]; p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = COLORS["PRIMARY"]; p.alignment = PP_ALIGN.CENTER

    # 5. 0% / 100% 끝 레이블
    end_y = label_y + int(label_h) + int(Inches(0.1))
    for pct_txt, tx in [("0%", bar_x - int(Inches(0.1))), ("100%", bar_x + bar_w - int(Inches(0.3)))]:
        etb = slide.shapes.add_textbox(tx, end_y, int(Inches(0.4)), int(Inches(0.3)))
        tf = etb.text_frame
        p = tf.paragraphs[0]; p.text = pct_txt
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(9)
        p.font.color.rgb = COLORS["GRAY"]; p.alignment = PP_ALIGN.CENTER

    # 6. 제목
    if title:
        title_y = end_y + int(Inches(0.35))
        # 경계 체크
        if title_y + int(title_h) > SLIDE_BOTTOM:
            title_y = int(SLIDE_BOTTOM) - int(title_h)
        ttb = slide.shapes.add_textbox(int(bx), title_y, int(bw), int(title_h))
        tf = ttb.text_frame
        p = tf.paragraphs[0]; p.text = title
        p.font.name = FONTS["BODY_TITLE"]; p.font.size = Pt(14); p.font.bold = True
        p.font.color.rgb = COLORS["DARK_GRAY"]; p.alignment = PP_ALIGN.CENTER


