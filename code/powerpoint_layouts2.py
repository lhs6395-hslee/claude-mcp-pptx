# -*- coding: utf-8 -*-
from powerpoint_utils import *

# 16. Full Image (풀와이드 이미지/다이어그램)
def render_full_image(slide, data):
    """이미지/다이어그램이 슬라이드 본문 전체를 차지하는 레이아웃

    data.data.data:
        image_path: 이미지 파일 경로
        search_q: architecture/ 폴더 검색어 (image_path 없을 때)
        caption: 하단 캡션 (선택)
    """
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})

    caption = content.get('caption', '')
    caption_h = Inches(0.45) if caption else 0
    caption_gap = Inches(0.1) if caption else 0
    img_h = bh - caption_h - caption_gap

    # 이미지 로드 시도
    img_loaded = False
    image_path = content.get('image_path', '')

    # 1순위: 직접 경로
    if image_path and os.path.exists(image_path):
        img_loaded = _place_image_centered(slide, image_path, bx, by, bw, img_h)

    # 2순위: architecture/ 폴더 검색
    if not img_loaded:
        search_q = content.get('search_q', '')
        if search_q:
            img_file = os.path.join('architecture', search_q.replace(' ', '_') + '.png')
            if os.path.exists(img_file):
                img_loaded = _place_image_centered(slide, img_file, bx, by, bw, img_h)

    # 3순위: screenshots/ 폴더 검색
    if not img_loaded:
        search_q = content.get('search_q', '')
        if search_q:
            img_file = os.path.join('screenshots', search_q.replace(' ', '_') + '.png')
            if os.path.exists(img_file):
                img_loaded = _place_image_centered(slide, img_file, bx, by, bw, img_h)

    # 폴백: 회색 박스
    if not img_loaded:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, img_h)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = COLORS["BG_BOX"]
        placeholder.line.color.rgb = COLORS["BORDER"]
        tf = placeholder.text_frame; tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"Image: {content.get('image_path', '') or content.get('search_q', 'N/A')}"
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(14)
        p.font.color.rgb = COLORS["GRAY"]; p.alignment = PP_ALIGN.CENTER

    # 캡션
    if caption:
        cap_y = by + img_h + caption_gap
        tb = slide.shapes.add_textbox(bx, cap_y, bw, caption_h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = caption
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(12)
        p.font.color.rgb = COLORS["GRAY"]; p.alignment = PP_ALIGN.CENTER


def _place_image_centered(slide, image_path, area_x, area_y, area_w, area_h):
    """이미지를 영역 내 중앙에 비율 유지하며 배치 (공통 유틸리티)"""
    try:
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            orig_w, orig_h = img.size
        aspect = orig_w / orig_h
        aw, ah = int(area_w), int(area_h)
        if aw / aspect <= ah:
            fw, fh = aw, int(aw / aspect)
        else:
            fh, fw = ah, int(ah * aspect)
        cx = int(area_x) + (aw - fw) // 2
        cy = int(area_y) + (ah - fh) // 2
        slide.shapes.add_picture(image_path, cx, cy, width=fw, height=fh)
        return True
    except ImportError:
        slide.shapes.add_picture(image_path, int(area_x), int(area_y),
                                 width=int(area_w), height=int(area_h))
        return True
    except Exception as e:
        print(f"   ⚠️ [이미지 로드 실패] {str(e)[:50]}")
        return False


# 17. Before/After (전/후 비교)
def render_before_after(slide, data):
    """Before/After 비교 레이아웃

    좌측: Before (회색/빨강 톤) / 우측: After (파랑/녹색 톤)
    중앙에 화살표

    data.data:
        before_title, before_body: Before 패널
        after_title, after_body: After 패널
    """
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    # before_after는 wrapper 레벨에서 직접 읽음 (예외)
    content = wrapper

    arrow_gap = Inches(0.8)
    w_half = (bw - arrow_gap) / 2
    label_h = Inches(0.55)
    body_gap = Inches(0.1)

    # ── Before 패널 (좌측) ──
    before_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, w_half, label_h)
    before_label.fill.solid()
    before_label.fill.fore_color.rgb = COLORS["SEM_RED"]
    before_label.line.color.rgb = COLORS["SEM_RED"]
    tf = before_label.text_frame; tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = content.get('before_title', 'Before')
    p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(18); p.font.color.rgb = COLORS["BG_WHITE"]
    p.alignment = PP_ALIGN.CENTER

    before_body_y = by + label_h + body_gap
    before_body_h = bh - label_h - body_gap
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, before_body_y, w_half, before_body_h)
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLORS["SEM_RED_BG"]
    before_box.line.color.rgb = COLORS["SEM_RED"]
    before_box.line.width = Pt(1.5)

    tf = before_box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    _before_lines = str(content.get('before_body', '')).split('\n')
    _before_is_list = len(_before_lines) > 1
    import re as _re
    for i, line in enumerate(_before_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.strip()
        is_numbered = bool(_re.match(r'^\d+[.):\s]\s', stripped))
        if _before_is_list and stripped and not stripped.startswith('•') and not is_numbered:
            p.text = "• " + line
        else:
            p.text = line
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(14)
        p.font.color.rgb = COLORS["SEM_RED_TEXT"]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6); p.line_spacing = 1.3

    # ── After 패널 (우측) ──
    after_x = bx + w_half + arrow_gap

    after_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, after_x, by, w_half, label_h)
    after_label.fill.solid()
    after_label.fill.fore_color.rgb = COLORS["SEM_GREEN"]
    after_label.line.color.rgb = COLORS["SEM_GREEN"]
    tf = after_label.text_frame; tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = content.get('after_title', 'After')
    p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
    p.font.size = Pt(18); p.font.color.rgb = COLORS["BG_WHITE"]
    p.alignment = PP_ALIGN.CENTER

    after_body_y = by + label_h + body_gap
    after_body_h = bh - label_h - body_gap
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, after_x, after_body_y, w_half, after_body_h)
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLORS["SEM_GREEN_BG"]
    after_box.line.color.rgb = COLORS["SEM_GREEN"]
    after_box.line.width = Pt(1.5)

    tf = after_box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    _after_lines = str(content.get('after_body', '')).split('\n')
    _after_is_list = len(_after_lines) > 1
    for i, line in enumerate(_after_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.strip()
        is_numbered = bool(_re.match(r'^\d+[.):\s]\s', stripped))
        if _after_is_list and stripped and not stripped.startswith('•') and not is_numbered:
            p.text = "• " + line
        else:
            p.text = line
        p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(14)
        p.font.color.rgb = COLORS["SEM_GREEN_TEXT"]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6); p.line_spacing = 1.3

    # ── 중앙 화살표 ──
    arrow_w = Inches(1.0); arrow_h_size = Inches(1.0)
    arrow_x = bx + w_half + (arrow_gap - arrow_w) / 2
    arrow_y = by + (bh / 2) - (arrow_h_size / 2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, arrow_w, arrow_h_size)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["PRIMARY"]
    arrow.line.color.rgb = COLORS["PRIMARY"]


# 18. Icon Grid (6~9 아이콘 그리드)
def render_icon_grid(slide, data):
    """아이콘 + 제목 + 설명 그리드 (6~9개 아이템)

    자동 레이아웃: 3열 x N행 (아이템 수에 따라)

    data.data.data.items: [
        {"icon": "kubernetes", "title": "제목", "desc": "설명"},
        ...
    ]
    """
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})

    items = content.get('items', [])
    if not items:
        return

    # 그리드 계산: 항상 3열
    n_cols = 3
    n_rows = (len(items) + n_cols - 1) // n_cols  # ceil

    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    cell_w = (bw - gap_x * (n_cols - 1)) / n_cols
    cell_h = (bh - gap_y * (n_rows - 1)) / n_rows

    icon_size = Inches(0.55)
    text_left_margin = icon_size + Inches(0.2)

    for idx, item in enumerate(items):
        col = idx % n_cols
        row = idx // n_cols
        x = bx + col * (cell_w + gap_x)
        y = by + row * (cell_h + gap_y)

        # 셀 배경 박스
        cell_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
        cell_box.fill.solid()
        cell_box.fill.fore_color.rgb = COLORS["BG_WHITE"]
        cell_box.line.color.rgb = COLORS["BORDER"]
        cell_box.line.width = Pt(1.0)

        # 아이콘 (좌측 상단)
        icon_x = x + Inches(0.15)
        icon_y = y + (cell_h - icon_size) / 2
        draw_icon_search(slide, icon_x, icon_y, icon_size, item.get('icon', ''))

        # 텍스트 (아이콘 우측)
        text_x = x + text_left_margin
        text_w = cell_w - text_left_margin - Inches(0.1)
        tb = slide.shapes.add_textbox(text_x, y + Inches(0.1), text_w, cell_h - Inches(0.2))
        tf = tb.text_frame; tf.word_wrap = True; tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)

        # 제목
        p = tf.paragraphs[0]
        p.text = item.get('title', '')
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(14); p.font.color.rgb = COLORS["PRIMARY"]
        p.space_after = Pt(4)

        # 설명
        desc = item.get('desc', '')
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(11)
            p2.font.color.rgb = COLORS["GRAY"]
            p2.line_spacing = 1.2


# 19. Numbered List (번호형 세로 리스트)
def render_numbered_list(slide, data):
    """번호형 세로 스텝 리스트

    좌측 큰 번호 원형 + 우측 제목/설명

    data.data.data.items: [
        {"title": "항목 제목", "desc": "항목 설명"},
        ...
    ]
    """
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})

    items = content.get('items', [])
    if not items:
        return

    n = len(items)
    # 항목 수에 따른 동적 크기 조절
    if n <= 4:
        gap = Inches(0.15); badge_size = Inches(0.65); title_pt = 16; desc_pt = 13; badge_pt = 22
    elif n <= 5:
        gap = Inches(0.12); badge_size = Inches(0.55); title_pt = 15; desc_pt = 12; badge_pt = 20
    elif n <= 6:
        gap = Inches(0.10); badge_size = Inches(0.48); title_pt = 13; desc_pt = 11; badge_pt = 18
    else:
        gap = Inches(0.08); badge_size = Inches(0.42); title_pt = 12; desc_pt = 10; badge_pt = 16

    item_h = (bh - gap * (n - 1)) / n
    text_left = badge_size + Inches(0.25)

    for i, item in enumerate(items):
        y = by + i * (item_h + gap)

        # 배경 바 (연한 회색)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, bw, item_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["BG_BOX"] if i % 2 == 0 else COLORS["BG_WHITE"]
        bar.line.color.rgb = COLORS["BORDER"]
        bar.line.width = Pt(1.0)

        # 번호 배지
        badge_x = bx + Inches(0.15)
        badge_y = y + (item_h - badge_size) / 2
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS["PRIMARY"]
        badge.line.color.rgb = COLORS["PRIMARY"]

        tf_badge = badge.text_frame; tf_badge.clear()
        tf_badge.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_badge.paragraphs[0]
        p.text = str(i + 1)
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(badge_pt); p.font.color.rgb = COLORS["BG_WHITE"]
        p.alignment = PP_ALIGN.CENTER

        # 텍스트 (제목 + 설명)
        text_x = bx + text_left + Inches(0.1)
        text_w = bw - text_left - Inches(0.3)
        tb = slide.shapes.add_textbox(text_x, y + Inches(0.05), text_w, item_h - Inches(0.1))
        tf = tb.text_frame; tf.word_wrap = True; tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)

        # 제목
        p = tf.paragraphs[0]
        p.text = item.get('title', '')
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(title_pt); p.font.color.rgb = COLORS["DARK_GRAY"]
        p.space_after = Pt(2)

        # 설명
        desc = item.get('desc', '')
        if desc:
            for j, line in enumerate(desc.split('\n')):
                p2 = tf.add_paragraph()
                p2.text = line
                p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(desc_pt)
                p2.font.color.rgb = COLORS["GRAY"]
                p2.space_after = Pt(1); p2.line_spacing = 1.1


# 20. Stats Dashboard (KPI/대형 숫자 표시)
def render_stats_dashboard(slide, data):
    """KPI/대형 숫자 대시보드 레이아웃

    3~4개 메트릭을 큰 숫자로 강조 표시

    data.data.data.metrics: [
        {"value": "99.9", "unit": "%", "label": "가용성", "desc": "연간 SLA 기준"},
        ...
    ]
    """
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})

    metrics = content.get('metrics', [])
    if not metrics:
        return

    n = len(metrics)
    gap = Inches(0.25)
    card_w = (bw - gap * (n - 1)) / n

    # 색상 팔레트 (순환)
    accent_colors = [
        (COLORS["PRIMARY"], COLORS["SEM_BLUE_BG"]),
        (COLORS["SEM_GREEN"], COLORS["SEM_GREEN_BG"]),
        (COLORS["SEM_ORANGE"], COLORS["SEM_ORANGE_BG"]),
        (COLORS["SEM_RED"], COLORS["SEM_RED_BG"]),
    ]

    for i, metric in enumerate(metrics):
        x = bx + i * (card_w + gap)
        accent, bg = accent_colors[i % len(accent_colors)]

        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, by, card_w, bh)
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = accent
        card.line.width = Pt(2.0)

        # 레이아웃: 상단 65% 숫자, 하단 35% 라벨+설명
        number_h = bh * 0.55
        label_h = bh * 0.45

        # 큰 숫자 + 단위
        tb_num = slide.shapes.add_textbox(x, by + Inches(0.2), card_w, number_h)
        tf = tb_num.text_frame; tf.clear(); tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        value_text = str(metric.get('value', ''))
        unit_text = metric.get('unit', '')
        p.alignment = PP_ALIGN.CENTER

        # value와 unit을 별도 run으로 분리 (크기 차이)
        run_val = p.add_run()
        run_val.text = value_text
        run_val.font.name = FONTS["BODY_TITLE"]; run_val.font.bold = True
        run_val.font.size = Pt(44); run_val.font.color.rgb = accent

        if unit_text:
            run_unit = p.add_run()
            run_unit.text = unit_text
            run_unit.font.name = FONTS["BODY_TITLE"]; run_unit.font.bold = True
            run_unit.font.size = Pt(24); run_unit.font.color.rgb = accent

        # 라벨
        tb_label = slide.shapes.add_textbox(x, by + number_h, card_w, label_h)
        tf = tb_label.text_frame; tf.clear(); tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = metric.get('label', '')
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(16); p.font.color.rgb = COLORS["DARK_GRAY"]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        # 설명
        desc = metric.get('desc', '')
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(12)
            p2.font.color.rgb = COLORS["GRAY"]
            p2.alignment = PP_ALIGN.CENTER
            p2.line_spacing = 1.2


# 21. Quote Highlight (인용문/핵심 메시지 강조)
def render_quote_highlight(slide, data):
    wrapper = data.get('data', {})
    y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start)
    content = wrapper.get('data', {})

    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLORS["SEM_BLUE_BG"]
    bg.line.color.rgb = COLORS["PRIMARY"]; bg.line.width = Pt(2.0)

    bar_w = Inches(0.08)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx + Inches(0.4), by + Inches(0.5), bar_w, bh - Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLORS["PRIMARY"]; bar.line.color.rgb = COLORS["PRIMARY"]

    tb_mark = slide.shapes.add_textbox(bx + Inches(0.7), by + Inches(0.2), Inches(1.0), Inches(0.8))
    p = tb_mark.text_frame.paragraphs[0]; p.text = "\u201C"
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = COLORS["PRIMARY"]

    quote_x = bx + Inches(0.8); quote_w = bw - Inches(1.6); quote_h = bh * 0.6
    tb_quote = slide.shapes.add_textbox(quote_x, by + Inches(0.8), quote_w, quote_h)
    tf = tb_quote.text_frame; tf.word_wrap = True; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = content.get('quote', '')
    p.font.name = FONTS["BODY_TITLE"]; p.font.size = Pt(22)
    p.font.italic = True; p.font.color.rgb = COLORS["DARK_GRAY"]
    p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.4

    author = content.get('author', ''); role = content.get('role', '')
    if author:
        author_y = by + bh - Inches(0.9)
        tb_author = slide.shapes.add_textbox(quote_x, author_y, quote_w, Inches(0.7))
        tf = tb_author.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; p.text = f"— {author}"
        p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["PRIMARY"]
        if role:
            p2 = tf.add_paragraph(); p2.text = role
            p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(13); p2.font.color.rgb = COLORS["GRAY"]


# 22. Pros & Cons (장단점 비교)
def render_pros_cons(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    gap = Inches(0.3); w_half = (bw - gap) / 2

    subject = content.get('subject', ''); subject_h = Inches(0.55) if subject else 0
    if subject:
        tb = slide.shapes.add_textbox(bx, by, bw, subject_h); tf = tb.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = subject; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(20); p.font.color.rgb = COLORS["DARK_GRAY"]; p.alignment = PP_ALIGN.CENTER

    panel_y = by + subject_h + (Inches(0.1) if subject else 0); panel_h = bh - subject_h - (Inches(0.1) if subject else 0); label_h = Inches(0.5)

    # Pros (좌측 - 녹색)
    pros_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, panel_y, w_half, label_h)
    pros_label.fill.solid(); pros_label.fill.fore_color.rgb = COLORS["SEM_GREEN"]; pros_label.line.color.rgb = COLORS["SEM_GREEN"]
    tf = pros_label.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = "PROS"; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    pros_body_y = panel_y + label_h + Inches(0.1); pros_body_h = panel_h - label_h - Inches(0.1)
    pros_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, pros_body_y, w_half, pros_body_h)
    pros_box.fill.solid(); pros_box.fill.fore_color.rgb = COLORS["SEM_GREEN_BG"]; pros_box.line.color.rgb = COLORS["SEM_GREEN"]; pros_box.line.width = Pt(1.5)

    tb = slide.shapes.add_textbox(bx + Inches(0.2), pros_body_y + Inches(0.15), w_half - Inches(0.4), pros_body_h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(content.get('pros', [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2714  {item}"; p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(14); p.font.color.rgb = COLORS["SEM_GREEN_TEXT"]; p.space_after = Pt(8); p.line_spacing = 1.3

    # Cons (우측 - 빨강)
    cons_x = bx + w_half + gap
    cons_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cons_x, panel_y, w_half, label_h)
    cons_label.fill.solid(); cons_label.fill.fore_color.rgb = COLORS["SEM_RED"]; cons_label.line.color.rgb = COLORS["SEM_RED"]
    tf = cons_label.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = "CONS"; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    cons_body_y = panel_y + label_h + Inches(0.1); cons_body_h = panel_h - label_h - Inches(0.1)
    cons_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cons_x, cons_body_y, w_half, cons_body_h)
    cons_box.fill.solid(); cons_box.fill.fore_color.rgb = COLORS["SEM_RED_BG"]; cons_box.line.color.rgb = COLORS["SEM_RED"]; cons_box.line.width = Pt(1.5)

    tb = slide.shapes.add_textbox(cons_x + Inches(0.2), cons_body_y + Inches(0.15), w_half - Inches(0.4), cons_body_h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(content.get('cons', [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2718  {item}"; p.font.name = FONTS["BODY_TEXT"]; p.font.size = Pt(14); p.font.color.rgb = COLORS["SEM_RED_TEXT"]; p.space_after = Pt(8); p.line_spacing = 1.3


# 23. Do / Don't (가이드라인 레이아웃)
def render_do_dont(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    gap = Inches(0.3); w_half = (bw - gap) / 2; label_h = Inches(0.6)

    # DO 패널 (좌측 - 파랑)
    do_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, w_half, label_h)
    do_label.fill.solid(); do_label.fill.fore_color.rgb = COLORS["PRIMARY"]; do_label.line.color.rgb = COLORS["PRIMARY"]
    tf = do_label.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = "\u2714  DO — 이렇게 하세요"
    p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    do_items = content.get('do_items', [])
    if do_items:
        item_y = by + label_h + Inches(0.15); item_gap = Inches(0.1)
        item_h = (bh - label_h - Inches(0.15) - item_gap * (len(do_items) - 1)) / len(do_items)
        for i, item in enumerate(do_items):
            iy = item_y + i * (item_h + item_gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, iy, w_half, item_h)
            box.fill.solid(); box.fill.fore_color.rgb = COLORS["SEM_BLUE_BG"]; box.line.color.rgb = COLORS["PRIMARY"]; box.line.width = Pt(1.0)
            text = item if isinstance(item, str) else item.get('text', ''); detail = '' if isinstance(item, str) else item.get('detail', '')
            tb = slide.shapes.add_textbox(bx + Inches(0.2), iy + Inches(0.08), w_half - Inches(0.4), item_h - Inches(0.16))
            tf = tb.text_frame; tf.word_wrap = True; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.text = f"\u2714  {text}"; p.font.name = FONTS["BODY_TEXT"]; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLORS["PRIMARY"]; p.space_after = Pt(2)
            if detail:
                p2 = tf.add_paragraph(); p2.text = detail; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(11); p2.font.color.rgb = COLORS["GRAY"]

    # DON'T 패널 (우측 - 빨강)
    dont_x = bx + w_half + gap
    dont_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dont_x, by, w_half, label_h)
    dont_label.fill.solid(); dont_label.fill.fore_color.rgb = COLORS["SEM_RED"]; dont_label.line.color.rgb = COLORS["SEM_RED"]
    tf = dont_label.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = "\u2718  DON'T — 이렇게 하지 마세요"
    p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    dont_items = content.get('dont_items', [])
    if dont_items:
        item_y = by + label_h + Inches(0.15); item_gap = Inches(0.1)
        item_h = (bh - label_h - Inches(0.15) - item_gap * (len(dont_items) - 1)) / len(dont_items)
        for i, item in enumerate(dont_items):
            iy = item_y + i * (item_h + item_gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dont_x, iy, w_half, item_h)
            box.fill.solid(); box.fill.fore_color.rgb = COLORS["SEM_RED_BG"]; box.line.color.rgb = COLORS["SEM_RED"]; box.line.width = Pt(1.0)
            text = item if isinstance(item, str) else item.get('text', ''); detail = '' if isinstance(item, str) else item.get('detail', '')
            tb = slide.shapes.add_textbox(dont_x + Inches(0.2), iy + Inches(0.08), w_half - Inches(0.4), item_h - Inches(0.16))
            tf = tb.text_frame; tf.word_wrap = True; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.text = f"\u2718  {text}"; p.font.name = FONTS["BODY_TEXT"]; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLORS["SEM_RED"]; p.space_after = Pt(2)
            if detail:
                p2 = tf.add_paragraph(); p2.text = detail; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(11); p2.font.color.rgb = COLORS["GRAY"]


# 24. Split Text + Code (좌측 설명 + 우측 코드 블록)
def render_split_text_code(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    gap = Inches(0.3); w_left = (bw - gap) * 0.4; w_right = (bw - gap) * 0.6

    desc = content.get('description', '')
    bullets = content.get('bullets', [])

    # 동적 폰트 크기 (텍스트 오버플로우 방지)
    n_desc = sum(max(1, len(line) // 38 + 1) for line in desc.split('\n')) if desc else 0
    n_total = n_desc + len(bullets)
    if n_total > 10:
        desc_sz, bul_sz, spc = Pt(11), Pt(10), Pt(3)
    elif n_total > 7:
        desc_sz, bul_sz, spc = Pt(13), Pt(12), Pt(4)
    else:
        desc_sz, bul_sz, spc = Pt(15), Pt(14), Pt(6)

    tb = slide.shapes.add_textbox(bx, by, w_left, bh)
    tf = tb.text_frame; tf.word_wrap = True; tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 텍스트 오버플로우 방지
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1); tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)

    if desc:
        for i, line in enumerate(desc.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.name = FONTS["BODY_TEXT"]; p.font.size = desc_sz; p.font.color.rgb = COLORS["DARK_GRAY"]; p.space_after = spc; p.line_spacing = 1.3

    if bullets:
        if desc: p_gap = tf.add_paragraph(); p_gap.text = ""; p_gap.space_after = Pt(6)
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if (desc or i > 0) else tf.paragraphs[0]
            p.text = f"• {bullet}"; p.font.name = FONTS["BODY_TEXT"]; p.font.size = bul_sz; p.font.color.rgb = COLORS["BLACK"]; p.space_after = spc; p.line_spacing = 1.2

    code_x = bx + w_left + gap
    lang = content.get('lang') or _detect_lang(content.get('code', ''), content.get('code_title', ''))
    create_terminal_box(slide, code_x, by, w_right, bh, content.get('code_title', 'code'), content.get('code', ''), lang=lang)


# 25. Pyramid Hierarchy (피라미드 계층 구조)
def render_pyramid_hierarchy(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    levels = content.get('levels', [])
    if not levels: return
    n = len(levels); gap = Inches(0.08); level_h = (bh - gap * (n - 1)) / n; center_x = bx + bw / 2
    min_w = bw * 0.3; max_w = bw * 0.95

    for i, level in enumerate(levels):
        ratio = i / max(n - 1, 1); level_w = min_w + (max_w - min_w) * ratio
        level_x = center_x - level_w / 2; level_y = by + i * (level_h + gap)
        color_key = level.get('color', 'primary')
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(level_x), int(level_y), int(level_w), int(level_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c; shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]; p.text = level.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(2)
        desc = level.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(12); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER


# 26. Cycle Loop (순환형 프로세스)
def render_cycle_loop(slide, data):
    import math
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    steps = content.get('steps', [])
    if not steps: return
    n = len(steps); center_label = content.get('center_label', '')
    cx = int(bx) + int(bw) // 2; cy = int(by) + int(bh) // 2
    # 와이드스크린 대응: 타원형 반경 (가로 넓게, 세로 좁게)
    radius_x = int(bw) // 2 - Inches(1.2)
    radius_y = int(bh) // 2 - Inches(0.8)

    if center_label:
        center_size = Inches(1.6)
        center_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - int(center_size) // 2, cy - int(center_size) // 2, int(center_size), int(center_size))
        center_shape.fill.solid(); center_shape.fill.fore_color.rgb = COLORS["PRIMARY"]; center_shape.line.color.rgb = COLORS["PRIMARY"]
        tf = center_shape.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = center_label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    step_colors = [
        (COLORS["PRIMARY"], COLORS["SEM_BLUE_BG"]), (COLORS["SEM_GREEN"], COLORS["SEM_GREEN_BG"]),
        (COLORS["SEM_ORANGE"], COLORS["SEM_ORANGE_BG"]), (COLORS["SEM_RED"], COLORS["SEM_RED_BG"]),
        (RGBColor(30, 58, 138), RGBColor(239, 246, 255)), (RGBColor(4, 120, 87), RGBColor(236, 253, 245)),
        (RGBColor(194, 65, 12), RGBColor(255, 247, 237)), (RGBColor(185, 28, 28), RGBColor(254, 242, 242)),
    ]

    box_w = Inches(1.8); box_h = Inches(1.2)
    for i, step in enumerate(steps):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        sx = cx + int(radius_x * math.cos(angle)) - int(box_w) // 2
        sy = cy + int(radius_y * math.sin(angle)) - int(box_h) // 2
        accent, bg = step_colors[i % len(step_colors)]

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, int(box_w), int(box_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = bg; shp.line.color.rgb = accent; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]; p.text = step.get('label', ''); p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = accent; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(2)
        desc = step.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(9); p2.font.color.rgb = COLORS["GRAY"]; p2.alignment = PP_ALIGN.CENTER

        # 화살표 (현재 → 다음 단계 방향)
        next_i = (i + 1) % n; next_angle = -math.pi / 2 + (2 * math.pi * next_i / n)
        mid_angle = (angle + next_angle) / 2
        if next_i == 0 and i == n - 1: mid_angle = angle + (2 * math.pi / n) / 2
        arrow_rx = int(radius_x * 0.65); arrow_ry = int(radius_y * 0.65)
        arrow_x = cx + int(arrow_rx * math.cos(mid_angle)) - Inches(0.15)
        arrow_y = cy + int(arrow_ry * math.sin(mid_angle)) - Inches(0.15)
        arrow_size = Inches(0.3)
        arrow = slide.shapes.add_shape(MSO_SHAPE.OVAL, arrow_x, arrow_y, int(arrow_size), int(arrow_size))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = COLORS["PRIMARY"]; arrow.line.color.rgb = COLORS["PRIMARY"]
        tf = arrow.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = "\u27A4"; p.font.size = Pt(10); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER


# 27. Venn Diagram (벤 다이어그램)
def render_venn_diagram(slide, data):
    """좌측 3원 벤 다이어그램 + 우측 설명 카드 레이아웃"""
    import math
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    circles = content.get('circles', [])
    if not circles: return
    center_label = content.get('center_label', '')
    n = min(len(circles), 3)

    # 좌측 55% = 원 영역, 우측 45% = 설명 카드
    left_w = int(bw * 0.55)
    right_x = int(bx) + left_w + Inches(0.25)
    right_w = int(bw) - left_w - Inches(0.25)
    vcx = int(bx) + left_w // 2; vcy = int(by) + int(bh) // 2

    # 파스텔 fill + 진한 border/text 색상
    circle_styles = [
        (RGBColor(219, 234, 254), RGBColor(30, 64, 175)),    # blue
        (RGBColor(254, 226, 226), RGBColor(185, 28, 28)),    # red
        (RGBColor(220, 252, 231), RGBColor(22, 101, 52)),    # green
    ]
    _venn_color_map = {'blue': 0, 'green': 2, 'orange': 1, 'red': 1}

    side = min(left_w, int(bh))
    circle_d = int(side * 0.55)
    offset = int(circle_d * 0.28)

    from pptx.oxml.ns import qn as _qn
    # 원 3개 배치 (정삼각형 꼭짓점) + 내부 라벨
    for i in range(n):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        ccx = vcx + int(offset * math.cos(angle)); ccy = vcy + int(offset * math.sin(angle))
        ox = ccx - circle_d // 2; oy = ccy - circle_d // 2

        ci = _venn_color_map.get(circles[i].get('color', ''), i % 3)
        fill_c, border_c = circle_styles[ci]

        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox, oy, circle_d, circle_d)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c
        try:
            sf = shp.fill._fill.find(_qn('a:solidFill'))
            if sf is not None:
                sc = sf.find(_qn('a:srgbClr'))
                if sc is not None:
                    a_el = sc.makeelement(_qn('a:alpha'), {}); a_el.set('val', '45000'); sc.append(a_el)
        except Exception: pass
        shp.line.color.rgb = border_c; shp.line.width = Pt(2.0)

        # 원 내부 라벨 (중심에서 바깥쪽으로 오프셋)
        label = circles[i].get('label', '')
        lbl_r = int(circle_d * 0.18)
        lx = ccx + int(lbl_r * math.cos(angle)); ly = ccy + int(lbl_r * math.sin(angle))
        lbl_w = Inches(1.8); lbl_h = Inches(0.45)
        lbl_shp = slide.shapes.add_textbox(lx - int(lbl_w) // 2, ly - int(lbl_h) // 2, int(lbl_w), int(lbl_h))
        tf = lbl_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(12); p.font.color.rgb = border_c; p.alignment = PP_ALIGN.CENTER

    # 중앙 교집합 라벨
    if center_label:
        cl_w = Inches(1.5); cl_h = Inches(0.55)
        c_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vcx - int(cl_w) // 2, vcy - int(cl_h) // 2, int(cl_w), int(cl_h))
        c_shp.fill.solid(); c_shp.fill.fore_color.rgb = COLORS["PRIMARY"]
        c_shp.line.color.rgb = COLORS["PRIMARY"]; c_shp.line.width = Pt(2.0)
        tf = c_shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = center_label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True
        p.font.size = Pt(13); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER

    # 우측 설명 카드 (컬러 왼쪽 액센트 바 + 흰색 카드)
    card_gap = Inches(0.15)
    card_h = (int(bh) - int(card_gap) * (n - 1)) // n
    bar_w = Inches(0.06)
    for i in range(n):
        ci = _venn_color_map.get(circles[i].get('color', ''), i % 3)
        _, border_c = circle_styles[ci]
        card_y = int(by) + (card_h + int(card_gap)) * i

        # 컬러 액센트 바 (좌측)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, card_y, int(bar_w), card_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = border_c; bar.line.fill.background()

        # 카드 본체
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       right_x + int(bar_w), card_y, right_w - int(bar_w), card_h)
        card.fill.solid(); card.fill.fore_color.rgb = COLORS["BG_WHITE"]
        card.line.color.rgb = RGBColor(229, 231, 235); card.line.width = Pt(1.0)
        tf = card.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)

        p = tf.paragraphs[0]; p.text = circles[i].get('label', ''); p.font.name = FONTS["BODY_TITLE"]
        p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = border_c
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(4)
        desc = circles[i].get('desc', '')
        if desc:
            for line in desc.split('\n'):
                p2 = tf.add_paragraph(); p2.text = f"• {line}"; p2.font.name = FONTS["BODY_TEXT"]
                p2.font.size = Pt(10); p2.font.color.rgb = COLORS["GRAY"]
                p2.alignment = PP_ALIGN.LEFT; p2.space_after = Pt(1)


# 28. SWOT Matrix (SWOT 분석 매트릭스)
def render_swot_matrix(slide, data):
    """2×2 그리드 + 중앙 축 라벨 (S/W/O/T 또는 커스텀)

    data.data.data:
        quadrants: [
            {"label": "S", "title": "Strengths", "items": ["기술력", "경험"], "color": "blue"},
            {"label": "W", "title": "Weaknesses", "items": ["인력부족"], "color": "red"},
            {"label": "O", "title": "Opportunities", "items": ["시장확대"], "color": "green"},
            {"label": "T", "title": "Threats", "items": ["경쟁심화"], "color": "orange"}
        ]
    """
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    quadrants = content.get('quadrants', [])
    if len(quadrants) < 4: return

    gap = Inches(0.5)   # 중앙 라벨 영역
    cell_w = (bw - gap) / 2; cell_h = (bh - gap) / 2

    positions = [
        (bx, by),                           # 좌상 (S)
        (bx + cell_w + gap, by),             # 우상 (W)
        (bx, by + cell_h + gap),             # 좌하 (O)
        (bx + cell_w + gap, by + cell_h + gap),  # 우하 (T)
    ]

    default_colors = ['blue', 'red', 'green', 'orange']

    for i, q in enumerate(quadrants[:4]):
        qx, qy = positions[i]
        color_key = q.get('color', default_colors[i])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        # 사분면 박스
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(qx), int(qy), int(cell_w), int(cell_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c
        shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)

        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.15)

        # 제목 (예: "Strengths")
        title = q.get('title', '')
        p = tf.paragraphs[0]; p.text = title; p.font.name = FONTS["BODY_TITLE"]
        p.font.bold = True; p.font.size = Pt(15); p.font.color.rgb = line_c; p.space_after = Pt(6)

        # 항목 리스트
        items = q.get('items', [])
        for item in items:
            p2 = tf.add_paragraph(); p2.text = f"• {item}"; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(12); p2.font.color.rgb = text_c; p2.space_after = Pt(3); p2.line_spacing = 1.2

    # 중앙 라벨 (S/W/O/T)
    label_size = Inches(0.45)
    label_positions = [
        (bx + cell_w - label_size / 2 + gap / 2, by + cell_h - label_size / 2 + gap / 2),  # 중앙
    ]
    # 4개 라벨을 중앙 교차점에 배치
    cx = int(bx + cell_w + gap / 2); cy_mid = int(by + cell_h + gap / 2)
    labels_pos = [
        (cx - int(label_size) - Inches(0.02), cy_mid - int(label_size) - Inches(0.02)),  # S (좌상)
        (cx + Inches(0.02), cy_mid - int(label_size) - Inches(0.02)),                    # W (우상)
        (cx - int(label_size) - Inches(0.02), cy_mid + Inches(0.02)),                    # O (좌하)
        (cx + Inches(0.02), cy_mid + Inches(0.02)),                                       # T (우하)
    ]

    for i, q in enumerate(quadrants[:4]):
        lx, ly = labels_pos[i]
        color_key = q.get('color', default_colors[i])
        _, line_c, _ = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(lx), int(ly), int(label_size), int(label_size))
        lbl.fill.solid(); lbl.fill.fore_color.rgb = line_c; lbl.line.color.rgb = line_c
        tf = lbl.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = q.get('label', ''); p.font.name = FONTS["BODY_TITLE"]
        p.font.bold = True; p.font.size = Pt(18); p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER


# 29. Center Radial (중심 방사형 관계도)
def render_center_radial(slide, data):
    """중앙 원 + 4방향 화살표 + 코너 라벨/설명

    data.data.data:
        center: {"label": "핵심 전략", "desc": "디지털 트랜스포메이션"}
        directions: [
            {"label": "기술", "desc": "클라우드, AI, DevOps", "color": "blue"},
            {"label": "프로세스", "desc": "자동화, 표준화", "color": "green"},
            {"label": "인력", "desc": "역량강화, 교육", "color": "orange"},
            {"label": "문화", "desc": "혁신, 협업", "color": "red"}
        ]
    """
    import math
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    center = content.get('center', {}); directions = content.get('directions', [])
    if not directions: return

    cx = int(bx) + int(bw) // 2; cy = int(by) + int(bh) // 2

    # 중앙 노드 (ROUNDED_RECTANGLE — OVAL은 내접 텍스트영역이 좁아 단어 잘림 발생)
    center_w = int(int(bh) * 0.42); center_h = int(int(bh) * 0.30)
    center_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx - center_w // 2, cy - center_h // 2, center_w, center_h)
    center_shp.fill.solid(); center_shp.fill.fore_color.rgb = COLORS["PRIMARY"]
    center_shp.line.color.rgb = COLORS["PRIMARY"]; center_shp.line.width = Pt(3.0)
    tf = center_shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    c_label = center.get('label', '')
    p = tf.paragraphs[0]; p.text = c_label; p.font.name = FONTS["BODY_TITLE"]
    c_fs = Pt(11) if len(c_label) > 16 else Pt(13) if len(c_label) > 10 else Pt(15)
    p.font.bold = True; p.font.size = c_fs; p.font.color.rgb = COLORS["BG_WHITE"]; p.alignment = PP_ALIGN.CENTER
    c_desc = center.get('desc', '')
    if c_desc:
        p2 = tf.add_paragraph(); p2.text = c_desc; p2.font.name = FONTS["BODY_TEXT"]
        p2.font.size = Pt(9); p2.font.color.rgb = RGBColor(200, 215, 255); p2.alignment = PP_ALIGN.CENTER

    # 4방향: 상, 우, 하, 좌 — 슬라이드 경계 보장 + 균일 간격
    n = min(len(directions), 4)
    default_colors_r = ['blue', 'green', 'orange', 'red']
    card_w = Inches(2.5); card_h = Inches(1.1)
    cr_v = center_h // 2; cr_h = center_w // 2  # 상하/좌우 반경 다름
    # 동적 간격: 상하 가용 공간 기준 + 좌우는 조금 더 길게
    avail_v = int(bh) // 2 - cr_v - int(card_h)
    v_gap = max(Inches(0.15), avail_v - Inches(0.08))
    h_gap = v_gap + Inches(0.4)

    card_positions = [
        (cx, cy - cr_v - int(v_gap) - int(card_h) // 2, 'top'),
        (cx + cr_h + int(h_gap) + int(card_w) // 2, cy, 'right'),
        (cx, cy + cr_v + int(v_gap) + int(card_h) // 2, 'bottom'),
        (cx - cr_h - int(h_gap) - int(card_w) // 2, cy, 'left'),
    ]

    for i in range(n):
        color_key = directions[i].get('color', default_colors_r[i])
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        card_cx, card_cy, direction = card_positions[i]
        card_x = card_cx - int(card_w) // 2; card_y = card_cy - int(card_h) // 2

        # 연결선 (중앙 노드 테두리 → 카드 테두리)
        if direction == 'top':
            lx1, ly1 = cx, cy - cr_v
            lx2, ly2 = cx, card_cy + int(card_h) // 2
        elif direction == 'bottom':
            lx1, ly1 = cx, cy + cr_v
            lx2, ly2 = cx, card_cy - int(card_h) // 2
        elif direction == 'right':
            lx1, ly1 = cx + cr_h, cy
            lx2, ly2 = card_cx - int(card_w) // 2, cy
        else:  # left
            lx1, ly1 = cx - cr_h, cy
            lx2, ly2 = card_cx + int(card_w) // 2, cy

        connector = slide.shapes.add_connector(1, lx1, ly1, lx2, ly2)
        connector.line.color.rgb = line_c; connector.line.width = Pt(2.5)

        # 카드
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, int(card_w), int(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = fill_c
        card.line.color.rgb = line_c; card.line.width = Pt(2.0)
        tf = card.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]; p.text = directions[i].get('label', ''); p.font.name = FONTS["BODY_TITLE"]
        p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
        d_desc = directions[i].get('desc', '')
        if d_desc:
            p2 = tf.add_paragraph(); p2.text = d_desc; p2.font.name = FONTS["BODY_TEXT"]
            p2.font.size = Pt(11); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER


# 30. Funnel (퍼널 다이어그램)
def render_funnel(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    stages = content.get('stages', [])
    if not stages: return
    n = len(stages); gap = Inches(0.06); level_h = (bh - gap * (n - 1)) / n; center_x = bx + bw / 2
    max_w = bw * 0.95; min_w = bw * 0.25

    for i, stage in enumerate(stages):
        ratio = i / max(n - 1, 1); level_w = max_w - (max_w - min_w) * ratio
        level_x = center_x - level_w / 2; level_y = by + i * (level_h + gap)
        color_key = stage.get('color', 'primary')
        fill_c, line_c, text_c = _SEM_BOX_STYLES.get(color_key, _SEM_BOX_STYLES['primary'])

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(level_x), int(level_y), int(level_w), int(level_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c; shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)

        value = stage.get('value', '')
        label = stage.get('label', '')
        if value:
            p = tf.paragraphs[0]; p.text = f"{label}  {value}"; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(2)
        else:
            p = tf.paragraphs[0]; p.text = label; p.font.name = FONTS["BODY_TITLE"]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = line_c; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(2)
        desc = stage.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(11); p2.font.color.rgb = text_c; p2.alignment = PP_ALIGN.CENTER


# 31. Zigzag Timeline (지그재그 타임라인)
def render_zigzag_timeline(slide, data):
    wrapper = data.get('data', {}); y_start = draw_body_header_and_get_y(slide, wrapper.get('body_title'), wrapper.get('body_desc'))
    bx, by, bw, bh = calculate_dynamic_rect(y_start); content = wrapper.get('data', {})
    steps = content.get('steps', [])
    if not steps: return
    n = len(steps)

    # 카드 크기 및 레이아웃 계산
    card_w = min(Inches(2.0), int(bw / max(n, 4) * 1.3)); card_h = Inches(1.3)
    top_y = int(by); bottom_y = int(by + bh) - card_h
    mid_y = int(by) + int(bh) // 2
    col_step = int(bw) / max(n, 1)

    step_colors = [
        (COLORS["PRIMARY"], COLORS["SEM_BLUE_BG"]), (COLORS["SEM_GREEN"], COLORS["SEM_GREEN_BG"]),
        (COLORS["SEM_ORANGE"], COLORS["SEM_ORANGE_BG"]), (COLORS["SEM_RED"], COLORS["SEM_RED_BG"]),
        (RGBColor(30, 58, 138), RGBColor(239, 246, 255)), (RGBColor(4, 120, 87), RGBColor(236, 253, 245)),
        (RGBColor(194, 65, 12), RGBColor(255, 247, 237)), (RGBColor(185, 28, 28), RGBColor(254, 242, 242)),
    ]

    # 날짜 기반 파란색 자동 감지 (MM/DD 형식 날짜가 오늘 이전이면 PRIMARY 파란색 솔리드)
    import datetime as _dt; import re as _re
    _today = _dt.date.today()
    def _step_start_date(ds):
        m = _re.match(r'(\d{2})/(\d{2})', ds.strip())
        if m:
            try: return _dt.date(2026, int(m.group(1)), int(m.group(2)))
            except: pass
        return None

    # 중앙 가로선 (배경 타임라인)
    line_shp = slide.shapes.add_connector(1, int(bx) + Inches(0.2), mid_y, int(bx + bw) - Inches(0.2), mid_y)
    line_shp.line.color.rgb = COLORS["BORDER"]; line_shp.line.width = Pt(2.0)
    line_shp.line.dash_style = 2  # dash

    for i, step in enumerate(steps):
        accent, bg = step_colors[i % len(step_colors)]

        # MM/DD 날짜 기반 3단계 색상 (원본 p3 스타일)
        # 과거: SEM_BLUE_BG(EFF6FF) fill + PRIMARY(0043DA) border (시작된 단계)
        # 미래: BG_BOX(F8F9FA) fill + DARK_GRAY(505050) border (예정 단계)
        # 날짜 없음: 기존 인덱스 색상 유지
        _sd = _step_start_date(step.get('date', ''))
        _is_past = _sd is not None and _sd <= _today
        if _is_past:
            fill_c = COLORS["SEM_BLUE_BG"]; line_c = COLORS["PRIMARY"]  # EFF6FF fill, 0043DA border
            dt_clr = COLORS["PRIMARY"]; title_clr = COLORS["PRIMARY"]; desc_clr = COLORS["GRAY"]
        elif _sd is not None:
            fill_c = COLORS["BG_BOX"]; line_c = COLORS["GRAY"]  # F8F9FA fill, 505050 border
            dt_clr = COLORS["GRAY"]; title_clr = COLORS["GRAY"]; desc_clr = COLORS["GRAY"]
        else:
            fill_c = bg; line_c = accent  # 날짜 없는 경우: 인덱스 색상
            dt_clr = accent; title_clr = accent; desc_clr = COLORS["GRAY"]

        cx = int(bx) + int(col_step * i) + int(col_step - card_w) // 2
        is_top = (i % 2 == 0)
        card_y = top_y if is_top else bottom_y

        # 수직 연결선 (카드 → 중앙선)
        conn_x = cx + card_w // 2
        if is_top:
            conn_y1 = card_y + card_h; conn_y2 = mid_y
        else:
            conn_y1 = mid_y; conn_y2 = card_y
        connector = slide.shapes.add_connector(1, conn_x, conn_y1, conn_x, conn_y2)
        connector.line.color.rgb = line_c; connector.line.width = Pt(1.5)

        # 중앙선 위 마커 원
        marker_size = Inches(0.2)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, conn_x - int(marker_size) // 2, mid_y - int(marker_size) // 2, int(marker_size), int(marker_size))
        marker.fill.solid(); marker.fill.fore_color.rgb = line_c; marker.line.color.rgb = line_c

        # 카드
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_c; shp.line.color.rgb = line_c; shp.line.width = Pt(2.0)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)

        date = step.get('date', '')
        if date:
            p0 = tf.paragraphs[0]; p0.text = date; p0.font.name = FONTS["BODY_TEXT"]; p0.font.size = Pt(9); p0.font.color.rgb = dt_clr; p0.alignment = PP_ALIGN.CENTER; p0.space_after = Pt(2)
            p1 = tf.add_paragraph()
        else:
            p1 = tf.paragraphs[0]
        _title_parts = step.get('title', '').split('\n')
        p1.text = _title_parts[0]; p1.font.name = FONTS["BODY_TITLE"]; p1.font.bold = True; p1.font.size = Pt(12); p1.font.color.rgb = title_clr; p1.alignment = PP_ALIGN.CENTER; p1.space_after = Pt(2)
        for _tp in _title_parts[1:]:
            _pt = tf.add_paragraph(); _pt.text = _tp; _pt.font.name = FONTS["BODY_TITLE"]; _pt.font.bold = True; _pt.font.size = Pt(12); _pt.font.color.rgb = title_clr; _pt.alignment = PP_ALIGN.CENTER; _pt.space_after = Pt(2)
        desc = step.get('desc', '')
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.name = FONTS["BODY_TEXT"]; p2.font.size = Pt(9); p2.font.color.rgb = desc_clr; p2.alignment = PP_ALIGN.CENTER


