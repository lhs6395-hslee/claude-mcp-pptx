#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
mcp-pptx MCP 서버

기존 PowerPoint 생성 엔진(mcp-pptx)을 MCP 도구로 래핑.
Claude가 JSON 데이터를 직접 전달 → PPT 생성.

Transport: stdio
등록: claude mcp add --transport stdio pptx-generator -- python3 /path/to/server.py
"""

import os
import sys
import shutil
import zipfile
import re
import copy
from datetime import datetime

from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from mcp.server.fastmcp import FastMCP

# 경로 설정 — 자기 디렉터리 기준 (독립 실행 가능)
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

# 엔진 모듈 import 경로 추가
sys.path.insert(0, os.path.join(PROJECT_ROOT, "code"))

from powerpoint_cover import update_cover_slide
from powerpoint_toc import update_toc_slide
from powerpoint_content import render_slide_content, set_slide_title_area
from transform import flat_to_engine_format

mcp = FastMCP("pptx-generator")


# --- generate.py에서 복사한 유틸리티 ---

def _remove_all_sections(pptx_file):
    """PowerPoint 파일에서 모든 섹션 제거"""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_rm_sec")
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    try:
        with zipfile.ZipFile(pptx_file, "r") as z:
            z.extractall(temp_dir)
        xml_path = os.path.join(temp_dir, "ppt", "presentation.xml")
        if os.path.exists(xml_path):
            tree = etree.parse(xml_path)
            for elem in list(tree.getroot().iter()):
                if "sectionLst" in elem.tag:
                    elem.getparent().remove(elem)
            tree.write(xml_path, xml_declaration=True, encoding="UTF-8", standalone=True)
        with zipfile.ZipFile(pptx_file, "w", zipfile.ZIP_DEFLATED) as z:
            for root, _, files in os.walk(temp_dir):
                for f in files:
                    full = os.path.join(root, f)
                    z.write(full, os.path.relpath(full, temp_dir))
    except Exception:
        pass
    finally:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)


def _clone_slide(prs, slide_index):
    """슬라이드를 복제하여 마지막에 추가"""
    from pptx.oxml.xmlchemy import OxmlElement
    import copy as copy_module
    
    source_slide = prs.slides[slide_index]
    
    # 슬라이드 XML 복사
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 모든 도형 복사
    for shape in source_slide.shapes:
        # 플레이스홀더가 아닌 도형만 복사
        if not shape.is_placeholder:
            el = shape.element
            newel = copy_module.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide


def _move_slide(prs, old_index, new_index):
    """슬라이드 위치 이동"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


# --- MCP 도구 ---

def _validate_slide(slide, slide_number: int) -> list[str]:
    """슬라이드 텍스트박스 겹침 및 경계 초과 검증.
    위반 항목 문자열 리스트 반환. 정상이면 빈 리스트."""
    SLIDE_BOTTOM_EMU = int(Inches(7.5))
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    issues = []

    for i, s1 in enumerate(text_shapes):
        bot = s1.top + s1.height
        # 경계 초과
        if bot > SLIDE_BOTTOM_EMU:
            txt = s1.text_frame.text.strip()[:20].replace('\n', ' ')
            issues.append(f"slide {slide_number}: OVERFLOW bot={bot/914400:.2f}\" \"{txt}\"")
        # 겹침 (다른 텍스트박스와)
        for s2 in text_shapes[i+1:]:
            t1, b1, l1, r1 = s1.top, s1.top+s1.height, s1.left, s1.left+s1.width
            t2, b2, l2, r2 = s2.top, s2.top+s2.height, s2.left, s2.left+s2.width
            if t1 < b2 and t2 < b1 and l1 < r2 and l2 < r1:
                txt1 = s1.text_frame.text.strip()[:15].replace('\n', ' ')
                txt2 = s2.text_frame.text.strip()[:15].replace('\n', ' ')
                issues.append(f"slide {slide_number}: OVERLAP \"{txt1}\" vs \"{txt2}\"")
    return issues


def _validate_toc(prs) -> list[str]:
    """목차 항목 수 vs 실제 섹션 수 검증.
    섹션 번호를 슬라이드 제목에서 추출해 목차와 비교."""
    issues = []

    # 목차 슬라이드에서 항목 수 집계
    toc_items = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and "CONTENTS" in shape.text_frame.text:
                # 숫자 박스와 제목 박스 찾기
                for s in slide.shapes:
                    if not hasattr(s, "text_frame"):
                        continue
                    paras = [p.text.strip() for p in s.text_frame.paragraphs if p.text.strip()]
                    if paras and paras[0] not in ("CONTENTS",) and not paras[0].isdigit():
                        for t in paras:
                            if t and t != "CONTENTS" and not t.isdigit():
                                toc_items.append(t)
                break

    # 본문 슬라이드에서 섹션 번호 추출 (제목 패턴: "N-M. ...")
    section_nums = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                txt = shape.text_frame.text.strip()
                m = re.match(r'^(\d+)-\d+\.', txt)
                if m:
                    section_nums.add(int(m.group(1)))

    if section_nums:
        max_section = max(section_nums)
        if len(toc_items) < max_section:
            issues.append(
                f"TOC: {len(toc_items)}개 항목이지만 섹션은 {max_section}개 — "
                f"섹션 {len(toc_items)+1}~{max_section} 목차 누락"
            )

    return issues


@mcp.tool()
def create_presentation(
    cover_title: str,
    cover_subtitle: str,
    sections: list[dict],
    output_name: str | None = None,
) -> str:
    """Generate a PowerPoint presentation.

    Args:
        cover_title: Cover slide title. Use \\n for line breaks.
        cover_subtitle: Cover slide subtitle.
        sections: List of section objects. Each section has:
            - section_title (str): e.g. "1. Overview"
            - slides (list): Each slide has:
                - layout (str): Layout name (e.g. "3_cards", "zigzag_timeline")
                - title (str): Slide header left
                - description (str): Slide header right (optional)
                - body_title (str): Body section title (optional)
                - body_desc (str): Body section description (optional)
                - content (dict): Layout-specific data in FLAT format

        Available layouts: bento_grid, 3_cards, grid_2x2, quad_matrix,
            timeline_steps, process_arrow, phased_columns, architecture_wide,
            image_left, comparison_vs, key_metric, challenge_solution,
            detail_image, comparison_table, detail_sections, table_callout,
            full_image, before_after, icon_grid, numbered_list,
            stats_dashboard, quote_highlight, pros_cons, do_dont,
            split_text_code, pyramid_hierarchy, cycle_loop, venn_diagram,
            swot_matrix, center_radial, funnel, zigzag_timeline,
            fishbone_cause_effect, org_chart, temple_pillars, infinity_loop,
            speedometer_gauge, mind_map, checklist_2col, kanban_board,
            exec_summary

        output_name: Output filename without extension. Defaults to timestamp.

    Returns:
        Absolute path to the generated .pptx file.
    """
    # 엔진의 상대 경로(icons/, template/ 등)를 위해 작업 디렉터리 변경
    original_cwd = os.getcwd()
    os.chdir(PROJECT_ROOT)

    # 엔진의 print()가 MCP stdout을 오염시키지 않도록 stderr로 리다이렉트
    original_stdout = sys.stdout
    sys.stdout = sys.stderr

    try:
        TEMPLATE_FILE = "template/2025_PPT_Template_FINAL.pptx"
        if not os.path.exists(TEMPLATE_FILE):
            return f"ERROR: Template not found: {os.path.join(PROJECT_ROOT, TEMPLATE_FILE)}"

        if not output_name:
            output_name = f"mcp_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        OUTPUT_FILE = f"results/{output_name}.pptx"

        os.makedirs("results", exist_ok=True)

        IDX_COVER = 0
        IDX_TOC = 1
        IDX_BODY_SRC = 7

        # 기존 파일이 있으면 삭제 후 새로 생성 (중복 슬라이드 방지)
        if os.path.exists(OUTPUT_FILE):
            os.remove(OUTPUT_FILE)
            print(f"[NEW MODE] Removed existing file, creating fresh: {OUTPUT_FILE}", file=sys.stderr)
        else:
            print(f"[NEW MODE] Creating new file: {OUTPUT_FILE}", file=sys.stderr)
        shutil.copy(TEMPLATE_FILE, OUTPUT_FILE)
        _remove_all_sections(OUTPUT_FILE)
        prs = Presentation(OUTPUT_FILE)
        is_append_mode = False

        if len(prs.slides) <= IDX_BODY_SRC:
            return "ERROR: Template has insufficient slides"

        keeper_ids = []

        if is_append_mode:
            # 누적 모드: 기존 슬라이드 모두 유지
            for slide in prs.slides:
                keeper_ids.append(slide.slide_id)
            
            # 엔딩 슬라이드 찾기 (마지막 슬라이드)
            ending = prs.slides[len(prs.slides) - 1]
            ending_id = ending.slide_id
            
            # 새 슬라이드는 엔딩 바로 앞에 삽입
            insert_idx = len(prs.slides) - 1
            
            # 목차 업데이트 (기존 섹션 + 새 섹션 병합)
            # 기존 목차 항목 추출
            existing_toc = []
            toc_slide_indices = []
            
            # 모든 목차 슬라이드 찾기 (2번째부터 시작, CONTENTS 텍스트 포함)
            for i, slide in enumerate(prs.slides):
                if i == 0:  # 표지 제외
                    continue
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        if "CONTENTS" in shape.text_frame.text:
                            toc_slide_indices.append(i)
                            break
            
            # 기존 목차에서 항목 추출 (순서 보존)
            existing_toc = []
            existing_toc_set = set()  # 중복 체크용
            for idx in toc_slide_indices:
                # 제목 박스 찾기: 숫자가 아닌 텍스트를 가진 다중 문단 박스
                candidates = []
                for shape in prs.slides[idx].shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    paras = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    if paras and paras[0] != "CONTENTS" and not paras[0].isdigit():
                        candidates.append((shape.left, paras))
                # 왼쪽에서 오른쪽 순으로 정렬 후 제목 박스(숫자 아닌 것) 선택
                candidates.sort(key=lambda x: x[0])
                for _, paras in candidates:
                    for text in paras:
                        if text and text != "CONTENTS" and not text.isdigit():
                            if text not in existing_toc_set:
                                existing_toc.append(text)
                                existing_toc_set.add(text)
            
            # 새 섹션 제목 추가 (슬라이드가 있는 섹션만)
            new_toc = [re.sub(r"^[\d\.]+\s*", "", s["section_title"]) 
                       for s in sections if s.get("slides")]
            
            # 새 섹션 중 기존에 없는 것만 추가 (중복 방지)
            for item in new_toc:
                if item not in existing_toc_set:
                    existing_toc.append(item)
            
            # 최종 목차
            combined_toc = existing_toc
            
            # 목차 페이지 업데이트 (5개씩 분할)
            MAX_ITEMS_PER_PAGE = 5
            toc_pages_needed = (len(combined_toc) + MAX_ITEMS_PER_PAGE - 1) // MAX_ITEMS_PER_PAGE
            
            # 기존 목차 페이지 업데이트
            for page_idx in range(len(toc_slide_indices)):
                start = page_idx * MAX_ITEMS_PER_PAGE
                end = min(start + MAX_ITEMS_PER_PAGE, len(combined_toc))
                start_number = start + 1  # 1-based numbering
                update_toc_slide(prs.slides[toc_slide_indices[page_idx]], combined_toc[start:end], start_number)
            
            # 추가 목차 페이지 필요 시 생성
            if len(toc_slide_indices) > 0:
                first_toc_idx = toc_slide_indices[0]
            else:
                first_toc_idx = IDX_TOC
                
            for page_idx in range(len(toc_slide_indices), toc_pages_needed):
                # 첫 번째 목차 슬라이드를 복제
                new_toc = _clone_slide(prs, first_toc_idx)
                
                # 새 목차를 표지 다음에 삽입
                _move_slide(prs, len(prs.slides) - 1, 1 + page_idx)
                keeper_ids.append(new_toc.slide_id)
                
                # 목차 항목 업데이트
                start = page_idx * MAX_ITEMS_PER_PAGE
                end = min(start + MAX_ITEMS_PER_PAGE, len(combined_toc))
                start_number = start + 1  # 1-based numbering
                update_toc_slide(prs.slides[1 + page_idx], combined_toc[start:end], start_number)
                
                insert_idx += 1  # 본문 삽입 위치 조정
            
            # 표지는 그대로 유지
            # body_layout은 기존 본문 슬라이드에서 가져오기
            body_found = False
            for i in range(2, len(prs.slides)):
                # 목차가 아닌 첫 본문 슬라이드 찾기
                is_toc = False
                for shape in prs.slides[i].shapes:
                    if hasattr(shape, "text_frame") and "CONTENTS" in shape.text_frame.text:
                        is_toc = True
                        break
                if not is_toc:
                    body_layout = prs.slides[i].slide_layout
                    body_found = True
                    break
            
            if not body_found:
                return "ERROR: Cannot append - no body slides found"
        else:
            # 새 파일 모드: 표지, 목차 생성
            # 표지
            cover = prs.slides[IDX_COVER]
            keeper_ids.append(cover.slide_id)
            update_cover_slide(cover, cover_title, cover_subtitle)

            # 목차 (5개씩 분할)
            clean_toc = [re.sub(r"^[\d\.]+\s*", "", s["section_title"]) for s in sections]
            MAX_ITEMS_PER_PAGE = 5
            toc_pages_needed = (len(clean_toc) + MAX_ITEMS_PER_PAGE - 1) // MAX_ITEMS_PER_PAGE
            
            print(f"[TOC] {len(clean_toc)} sections → {toc_pages_needed} TOC pages needed", file=sys.stderr)
            
            # 첫 번째 목차
            toc = prs.slides[IDX_TOC]
            keeper_ids.append(toc.slide_id)
            update_toc_slide(toc, clean_toc[:MAX_ITEMS_PER_PAGE])
            print(f"[TOC] Page 1: items 1-{min(MAX_ITEMS_PER_PAGE, len(clean_toc))}", file=sys.stderr)
            
            # 추가 목차 페이지 생성 (6개 이상인 경우)
            if toc_pages_needed > 1:
                for page_idx in range(1, toc_pages_needed):
                    # 첫 번째 목차 슬라이드를 복제
                    new_toc_slide = _clone_slide(prs, IDX_TOC)
                    
                    # 새 목차를 올바른 위치로 이동 (표지 다음)
                    _move_slide(prs, len(prs.slides) - 1, 1 + page_idx)
                    keeper_ids.append(prs.slides[1 + page_idx].slide_id)
                    
                    # 목차 항목 업데이트
                    start = page_idx * MAX_ITEMS_PER_PAGE
                    end = min(start + MAX_ITEMS_PER_PAGE, len(clean_toc))
                    start_number = start + 1  # 1-based numbering
                    update_toc_slide(prs.slides[1 + page_idx], clean_toc[start:end], start_number)
                    print(f"[TOC] Page {page_idx + 1}: items {start + 1}-{end}", file=sys.stderr)

            # 엔딩 슬라이드 보존
            ending = prs.slides[len(prs.slides) - 1]
            ending_id = ending.slide_id
            keeper_ids.append(ending_id)
            
            insert_idx = 1 + toc_pages_needed  # 표지 + 목차 페이지들 다음
            body_layout = prs.slides[IDX_BODY_SRC].slide_layout
            print(f"[TOC] Body slides will start at index {insert_idx}", file=sys.stderr)
        # 본문 생성
        slide_count = 0

        for section in sections:
            for slide_flat in section.get("slides", []):
                slide = prs.slides.add_slide(body_layout)
                _move_slide(prs, len(prs.slides) - 1, insert_idx)

                # flat → engine format 변환
                slide_data = flat_to_engine_format(slide_flat)

                set_slide_title_area(
                    slide, slide_data.get("t", ""), slide_data.get("d", "")
                )
                render_slide_content(
                    slide, slide_data.get("l", "bento_grid"), slide_data
                )

                keeper_ids.append(slide.slide_id)
                insert_idx += 1
                slide_count += 1

        if not is_append_mode:
            # 새 파일 모드: 불필요한 슬라이드 삭제
            xml_slides = prs.slides._sldIdLst
            for i in range(len(prs.slides) - 1, -1, -1):
                if prs.slides[i].slide_id not in keeper_ids:
                    rId = xml_slides[i].rId
                    prs.part.drop_rel(rId)
                    del xml_slides[i]

        # 엔딩 슬라이드를 마지막으로 이동
        for i, s in enumerate(prs.slides):
            if s.slide_id == ending_id:
                _move_slide(prs, i, len(prs.slides) - 1)
                break

        # 저장
        prs.save(OUTPUT_FILE)
        abs_path = os.path.abspath(OUTPUT_FILE)
        total = len(prs.slides)

        # 렌더링된 슬라이드 검증
        prs_check = Presentation(OUTPUT_FILE)
        all_issues = []
        # 새로 추가된 슬라이드만 검증 (insert_idx 이전 ~ 엔딩 직전)
        check_start = insert_idx - slide_count
        check_end = insert_idx
        for ci in range(check_start, check_end):
            if 0 <= ci < len(prs_check.slides):
                all_issues.extend(_validate_slide(prs_check.slides[ci], ci + 1))
        # 목차 검증
        all_issues.extend(_validate_toc(prs_check))

        mode_str = "appended to" if is_append_mode else "generated"
        result = f"{mode_str.capitalize()} {total} slides ({slide_count} body added): {abs_path}"
        if all_issues:
            result += "\n⚠️  VALIDATION ISSUES:\n" + "\n".join(f"  - {v}" for v in all_issues)
        else:
            result += f"\n✅ Validation passed ({slide_count} slides checked)"
        return result

    finally:
        sys.stdout = original_stdout
        os.chdir(original_cwd)


@mcp.tool()
def update_slide(
    output_name: str,
    slide_number: int,
    slide_data: dict,
) -> str:
    """Update a specific slide in an existing presentation.

    Args:
        output_name: Filename without extension (e.g. "layout_showcase_v2").
        slide_number: 1-based slide number to replace (e.g. 8 = 8th slide).
        slide_data: Slide definition with same format as create_presentation slides:
            - layout (str): Layout name
            - title (str): Slide header left
            - description (str): Slide header right (optional)
            - body_title (str): Body section title (optional)
            - body_desc (str): Body section description (optional)
            - content (dict): Layout-specific data in FLAT format

    Returns:
        Result message with updated slide info.
    """
    original_cwd = os.getcwd()
    os.chdir(PROJECT_ROOT)
    original_stdout = sys.stdout
    sys.stdout = sys.stderr

    try:
        OUTPUT_FILE = f"results/{output_name}.pptx"
        if not os.path.exists(OUTPUT_FILE):
            return f"ERROR: File not found: {OUTPUT_FILE}"

        prs = Presentation(OUTPUT_FILE)
        total = len(prs.slides)

        if slide_number < 1 or slide_number > total:
            return f"ERROR: slide_number {slide_number} out of range (1-{total})"

        idx = slide_number - 1  # 0-based
        target_slide = prs.slides[idx]

        # 기존 슬라이드의 content shape 모두 제거
        # 헤더 영역(top < 2.0인치 = 1828800 EMU) 이하의 shape만 제거
        HEADER_BOUNDARY = int(Inches(2.0))
        to_remove = [s for s in target_slide.shapes if s.top >= HEADER_BOUNDARY]
        for s in to_remove:
            s._element.getparent().remove(s._element)

        # 기존 슬라이드에 직접 렌더링
        data = flat_to_engine_format(slide_data)
        set_slide_title_area(target_slide, data.get("t", ""), data.get("d", ""))
        render_slide_content(target_slide, data.get("l", "bento_grid"), data)

        prs.save(OUTPUT_FILE)

        # 검증
        prs_check = Presentation(OUTPUT_FILE)
        issues = _validate_slide(prs_check.slides[idx], slide_number)
        issues.extend(_validate_toc(prs_check))
        result = f"Updated slide {slide_number} in {OUTPUT_FILE} (total: {total} slides)"
        if issues:
            result += "\n⚠️  VALIDATION ISSUES:\n" + "\n".join(f"  - {v}" for v in issues)
        else:
            result += "\n✅ Validation passed"
        return result

    finally:
        sys.stdout = original_stdout
        os.chdir(original_cwd)


@mcp.tool()
def delete_slide(
    output_name: str,
    slide_number: int,
) -> str:
    """Delete a specific slide from an existing presentation.

    Args:
        output_name: Filename without extension (e.g. "layout_showcase_v2").
        slide_number: 1-based slide number to delete (e.g. 8 = 8th slide).
                      Cover (1), TOC (2), and ending (last) slides can also be deleted.

    Returns:
        Result message with remaining slide count.
    """
    original_cwd = os.getcwd()
    os.chdir(PROJECT_ROOT)
    original_stdout = sys.stdout
    sys.stdout = sys.stderr

    try:
        OUTPUT_FILE = f"results/{output_name}.pptx"
        if not os.path.exists(OUTPUT_FILE):
            return f"ERROR: File not found: {OUTPUT_FILE}"

        prs = Presentation(OUTPUT_FILE)
        total = len(prs.slides)

        if slide_number < 1 or slide_number > total:
            return f"ERROR: slide_number {slide_number} out of range (1-{total})"

        idx = slide_number - 1  # 0-based
        xml_slides = prs.slides._sldIdLst
        elem = xml_slides[idx]
        rId = elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        xml_slides.remove(elem)
        prs.slides.part._rels._rels.pop(rId, None)

        prs.save(OUTPUT_FILE)
        remaining = len(prs.slides)
        return f"Deleted slide {slide_number} from {OUTPUT_FILE} (remaining: {remaining} slides)"

    finally:
        sys.stdout = original_stdout
        os.chdir(original_cwd)



@mcp.tool()
def insert_slide(
    output_name: str,
    after_slide_number: int,
    slide_data: dict,
) -> str:
    """Insert a new slide at a specific position in an existing presentation.

    Args:
        output_name: Filename without extension (e.g. "layout_showcase_v2").
        after_slide_number: Insert after this slide number (1-based).
                            e.g. 8 = insert as slide 9, pushing existing 9+ down.
                            Use 0 to insert as the first slide.
        slide_data: Slide definition with same format as create_presentation slides:
            - layout (str): Layout name
            - title (str): Slide header left
            - description (str): Slide header right (optional)
            - body_title (str): Body section title (optional)
            - body_desc (str): Body section description (optional)
            - content (dict): Layout-specific data in FLAT format

    Returns:
        Result message with new slide count.
    """
    original_cwd = os.getcwd()
    os.chdir(PROJECT_ROOT)
    original_stdout = sys.stdout
    sys.stdout = sys.stderr

    try:
        OUTPUT_FILE = f"results/{output_name}.pptx"
        if not os.path.exists(OUTPUT_FILE):
            return f"ERROR: File not found: {OUTPUT_FILE}"

        prs = Presentation(OUTPUT_FILE)
        total = len(prs.slides)

        if after_slide_number < 0 or after_slide_number > total:
            return f"ERROR: after_slide_number {after_slide_number} out of range (0-{total})"

        # 삽입 위치 (0-based): after_slide_number 바로 뒤
        insert_idx = after_slide_number  # 0-based index to insert at

        # 레이아웃은 인접 본문 슬라이드에서 가져오기
        ref_idx = min(insert_idx, total - 1)
        body_layout = prs.slides[ref_idx].slide_layout

        # 새 슬라이드를 마지막에 추가 후 원하는 위치로 이동
        new_slide = prs.slides.add_slide(body_layout)
        _move_slide(prs, len(prs.slides) - 1, insert_idx)

        # 콘텐츠 렌더링
        data = flat_to_engine_format(slide_data)
        set_slide_title_area(new_slide, data.get("t", ""), data.get("d", ""))
        render_slide_content(new_slide, data.get("l", "bento_grid"), data)

        prs.save(OUTPUT_FILE)
        new_total = len(prs.slides)

        # 검증
        prs_check = Presentation(OUTPUT_FILE)
        issues = _validate_slide(prs_check.slides[insert_idx], insert_idx + 1)
        issues.extend(_validate_toc(prs_check))
        result = f"Inserted slide at position {insert_idx + 1} in {OUTPUT_FILE} (total: {new_total} slides)"
        if issues:
            result += "\n⚠️  VALIDATION ISSUES:\n" + "\n".join(f"  - {v}" for v in issues)
        else:
            result += "\n✅ Validation passed"
        return result

    finally:
        sys.stdout = original_stdout
        os.chdir(original_cwd)



@mcp.tool()
def export_pdf(
    output_name: str,
    pdf_name: str | None = None,
) -> str:
    """Export an existing presentation to PDF using Microsoft PowerPoint (macOS).

    Args:
        output_name: Filename without extension of the source .pptx (e.g. "my_presentation").
        pdf_name: Output PDF filename without extension. Defaults to same as output_name.

    Returns:
        Absolute path to the generated .pdf file, or an error message.
    """
    import subprocess

    original_cwd = os.getcwd()
    os.chdir(PROJECT_ROOT)

    try:
        pptx_path = os.path.abspath(f"results/{output_name}.pptx")
        if not os.path.exists(pptx_path):
            return f"ERROR: File not found: {pptx_path}"

        if not pdf_name:
            pdf_name = output_name
        pdf_path = os.path.abspath(f"results/{pdf_name}.pdf")

        # AppleScript로 PowerPoint에서 PDF 내보내기
        script = f'''
tell application "Microsoft PowerPoint"
    set pptFile to POSIX file "{pptx_path}"
    set pdfFile to POSIX file "{pdf_path}"
    open pptFile
    set theDoc to active presentation
    save theDoc in pdfFile as save as PDF
    close theDoc saving no
end tell
'''
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=120
        )

        if result.returncode != 0:
            err = result.stderr.strip()
            return f"ERROR: AppleScript failed — {err}"

        if not os.path.exists(pdf_path):
            return f"ERROR: PDF not created at {pdf_path}"

        size_kb = os.path.getsize(pdf_path) // 1024
        return f"✅ PDF exported: {pdf_path} ({size_kb} KB)"

    finally:
        os.chdir(original_cwd)


if __name__ == "__main__":
    mcp.run(transport="stdio")
